"""
title: Präsentations-Generator (PowerPoint)
author: MTGE
version: 1.2.0
license: MIT
description: Erzeugt aus strukturierten Folien-Daten eine fertig gestaltete PowerPoint-Datei (16:9) und liefert einen Download-Link im Chat.
requirements: python-pptx>=0.6.23, requests
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import inspect
import io
import json
import math
import os
import re
import types
import uuid
from datetime import datetime
from typing import Any, Callable, Dict, List, Optional, Tuple

from pydantic import BaseModel, Field

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# =============================================================================
#  Farb-Hilfsfunktionen
# =============================================================================


def _clean_hex(value: str, fallback: str = "") -> str:
    """Normalisiert '#aabbcc' / 'aabbcc' / 'abc' auf 'AABBCC'."""
    if not value:
        return fallback
    v = str(value).strip().lstrip("#")
    if len(v) == 3:
        v = "".join(c * 2 for c in v)
    if not re.fullmatch(r"[0-9a-fA-F]{6}", v):
        return fallback
    return v.upper()


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(_clean_hex(value, "000000"))


def _to_tuple(value: str) -> Tuple[int, int, int]:
    v = _clean_hex(value, "000000")
    return int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)


def _mix(a: str, b: str, t: float) -> str:
    """Mischt Farbe a mit b (t = 0..1 Anteil von b). Ersetzt Transparenz sauber."""
    ar, ag, ab = _to_tuple(a)
    br, bg, bb = _to_tuple(b)
    t = max(0.0, min(1.0, t))
    return "%02X%02X%02X" % (
        round(ar + (br - ar) * t),
        round(ag + (bg - ag) * t),
        round(ab + (bb - ab) * t),
    )


def _relative_luminance(color: str) -> float:
    def chan(c: int) -> float:
        s = c / 255.0
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4

    r, g, b = _to_tuple(color)
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast(a: str, b: str) -> float:
    la, lb = _relative_luminance(a), _relative_luminance(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def _readable_on(bg: str, light: str = "FFFFFF", dark: str = "121212") -> str:
    """Wählt automatisch die besser lesbare Textfarbe für einen Hintergrund."""
    return light if _contrast(bg, light) >= _contrast(bg, dark) else dark


def _ensure_contrast(fg: str, bg: str, minimum: float = 3.4) -> str:
    """Hellt/dunkelt eine Akzentfarbe so weit auf, bis sie auf bg lesbar ist."""
    target = "FFFFFF" if _relative_luminance(bg) < 0.35 else "111111"
    out = fg
    for step in range(0, 11):
        if _contrast(out, bg) >= minimum:
            return out
        out = _mix(fg, target, step / 10.0)
    return out


# =============================================================================
#  Themes
# =============================================================================
#  deep      = dunkler Bühnen-Hintergrund (Titel / Kapitel / Abschluss)
#  light     = heller Inhalts-Hintergrund
#  surface   = Kartenfläche auf hellem Grund
#  ink/muted = Textfarben auf hellem Grund
#  accent*   = Akzente, series = Diagrammfarben

THEMES: Dict[str, Dict[str, Any]] = {
    "midnight": {
        "label": "Midnight Executive – Marine, souverän, klassisch Business",
        "deep": "0F1836", "deep2": "1B2550", "on_deep": "FFFFFF", "on_deep_mut": "B7C3E8",
        "light": "FFFFFF", "surface": "F2F4FA", "surface2": "E8ECF7",
        "ink": "111A33", "muted": "5A6784", "hair": "DCE1EE",
        "accent": "3E6DE0", "accent2": "E0A032", "accent3": "2AA391",
        "series": ["3E6DE0", "E0A032", "2AA391", "8B5CF6", "E0655A", "2E9BD6"],
    },
    "graphite": {
        "label": "Charcoal Minimal – reduziert, hoher Kontrast, ein scharfer Akzent",
        "deep": "1C2128", "deep2": "2A313A", "on_deep": "FFFFFF", "on_deep_mut": "B4BCC6",
        "light": "FFFFFF", "surface": "F4F5F7", "surface2": "E9ECEF",
        "ink": "1C2128", "muted": "5F6975", "hair": "E1E4E8",
        "accent": "D8492A", "accent2": "2F6F8F", "accent3": "7A8794",
        "series": ["D8492A", "2F6F8F", "7A8794", "C08A20", "3F7A5C", "6B5B95"],
    },
    "teal": {
        "label": "Teal Trust – frisch, technisch, vertrauenswürdig",
        "deep": "04353F", "deep2": "0A4C58", "on_deep": "FFFFFF", "on_deep_mut": "A9CBD2",
        "light": "FFFFFF", "surface": "EFF6F6", "surface2": "E0EDEE",
        "ink": "0B2A31", "muted": "4F6E76", "hair": "D6E4E6",
        "accent": "028090", "accent2": "01A183", "accent3": "D98A2B",
        "series": ["028090", "01A183", "D98A2B", "3E6DE0", "B8546B", "6E8E3F"],
    },
    "ocean": {
        "label": "Ocean Gradient – tiefes Blau, ruhig, seriös",
        "deep": "16234A", "deep2": "21295C", "on_deep": "FFFFFF", "on_deep_mut": "AFBBDA",
        "light": "FFFFFF", "surface": "F1F4F9", "surface2": "E4EAF3",
        "ink": "14203F", "muted": "566380", "hair": "DAE1EC",
        "accent": "065A82", "accent2": "1C7293", "accent3": "C98A2B",
        "series": ["065A82", "1C7293", "C98A2B", "4C9A8F", "8360A8", "C0584F"],
    },
    "forest": {
        "label": "Forest & Moss – natürlich, nachhaltig, geerdet",
        "deep": "16301C", "deep2": "234A2B", "on_deep": "FFFFFF", "on_deep_mut": "BBD3BE",
        "light": "FFFFFF", "surface": "F1F6F0", "surface2": "E3EDE1",
        "ink": "152218", "muted": "556B57", "hair": "DBE6D9",
        "accent": "2C5F2D", "accent2": "77A03F", "accent3": "C08A22",
        "series": ["2C5F2D", "77A03F", "C08A22", "3E7C8F", "7A5C9E", "B85042"],
    },
    "terracotta": {
        "label": "Warm Terracotta – warm, handwerklich, menschlich",
        "deep": "3B2320", "deep2": "543430", "on_deep": "FFFFFF", "on_deep_mut": "DCC3BC",
        "light": "FFFFFF", "surface": "F8F2EE", "surface2": "EFE4DD",
        "ink": "30201D", "muted": "6E5A54", "hair": "E7DCD5",
        "accent": "B85042", "accent2": "6E8C7C", "accent3": "C08A2B",
        "series": ["B85042", "6E8C7C", "C08A2B", "5E7A96", "8C6A9E", "5F7A5A"],
    },
    "coral": {
        "label": "Coral Energy – energiegeladen, jung, Marketing",
        "deep": "232C56", "deep2": "2F3C7E", "on_deep": "FFFFFF", "on_deep_mut": "C2C9E8",
        "light": "FFFFFF", "surface": "F5F6FB", "surface2": "E9ECF6",
        "ink": "1E2445", "muted": "5C6488", "hair": "DEE1EF",
        "accent": "E8474E", "accent2": "2F3C7E", "accent3": "D99A17",
        "series": ["E8474E", "2F3C7E", "D99A17", "2FA796", "8C6FD1", "5A7D9A"],
    },
    "berry": {
        "label": "Berry & Cream – edel, kulinarisch, weiblich-elegant",
        "deep": "3A1A28", "deep2": "55283A", "on_deep": "FFFFFF", "on_deep_mut": "D9BFC8",
        "light": "FFFFFF", "surface": "F8F1F2", "surface2": "EEE0E3",
        "ink": "2E141F", "muted": "6B4E58", "hair": "E8D9DC",
        "accent": "6D2E46", "accent2": "A26769", "accent3": "B07A45",
        "series": ["6D2E46", "A26769", "B07A45", "5E7A80", "7A8A4B", "B2545E"],
    },
    "cherry": {
        "label": "Cherry Bold – fast monochrom mit einem kräftigen Rot",
        "deep": "141414", "deep2": "242424", "on_deep": "FFFFFF", "on_deep_mut": "BDBDBD",
        "light": "FFFFFF", "surface": "F5F5F5", "surface2": "EAEAEA",
        "ink": "141414", "muted": "5F5F5F", "hair": "E3E3E3",
        "accent": "990011", "accent2": "2F3C7E", "accent3": "7A7A7A",
        "series": ["990011", "2F3C7E", "7A7A7A", "B4691F", "3F6F5B", "6B4E8C"],
    },
}

THEME_ALIASES = {
    "auto": "", "default": "", "standard": "",
    "business": "midnight", "corporate": "midnight", "seriös": "midnight",
    "blau": "midnight", "blue": "midnight", "navy": "midnight", "dunkel": "midnight",
    "minimal": "graphite", "minimalistisch": "graphite", "grau": "graphite",
    "monochrom": "cherry", "rot": "cherry", "red": "cherry", "bold": "cherry",
    "tech": "teal", "technik": "teal", "türkis": "teal", "startup": "teal",
    "wasser": "ocean", "finance": "ocean", "finanzen": "ocean",
    "grün": "forest", "green": "forest", "nachhaltig": "forest", "umwelt": "forest",
    "warm": "terracotta", "erdig": "terracotta", "orange": "terracotta",
    "marketing": "coral", "energie": "coral", "kreativ": "coral",
    "elegant": "berry", "premium": "berry", "lila": "berry",
}


def resolve_theme(name: str, fallback: str = "midnight") -> str:
    key = (name or "").strip().lower()
    key = THEME_ALIASES.get(key, key)
    if key in THEMES:
        return key
    fb = (fallback or "midnight").strip().lower()
    return fb if fb in THEMES else "midnight"


# =============================================================================
#  Layout-Konstanten (Zoll, 16:9)
# =============================================================================

SW, SH = 13.333, 7.5          # Foliengröße
MX = 0.90                     # Seitenrand links/rechts
CW = SW - 2 * MX              # nutzbare Breite = 11.533
TITLE_Y = 0.62
FOOTER_Y = 6.92
BODY_BOTTOM = 6.62


# =============================================================================
#  Slide-Normalisierung (verzeiht Modell-Fehler)
# =============================================================================

TYPE_ALIASES = {
    "title": "title", "titel": "title", "cover": "title", "deckblatt": "title", "start": "title",
    "section": "section", "abschnitt": "section", "kapitel": "section", "divider": "section",
    "trenner": "section", "chapter": "section",
    "bullets": "bullets", "bullet": "bullets", "list": "bullets", "liste": "bullets",
    "aufzählung": "bullets", "aufzaehlung": "bullets", "content": "bullets", "inhalt": "bullets",
    "text": "bullets", "punkte": "bullets",
    "agenda": "agenda", "überblick": "agenda", "ueberblick": "agenda", "toc": "agenda",
    "inhaltsverzeichnis": "agenda", "outline": "agenda",
    "cards": "cards", "karten": "cards", "kacheln": "cards", "grid": "cards", "tiles": "cards",
    "features": "cards", "boxes": "cards",
    "stats": "stats", "kpi": "stats", "kpis": "stats", "kennzahlen": "stats", "zahlen": "stats",
    "metrics": "stats", "numbers": "stats",
    "quote": "quote", "zitat": "quote", "testimonial": "quote", "statement": "quote",
    "compare": "compare", "vergleich": "compare", "two_column": "compare", "twocolumn": "compare",
    "columns": "compare", "spalten": "compare", "prosandcons": "compare", "vorher_nachher": "compare",
    "table": "table", "tabelle": "table", "matrix": "table",
    "timeline": "timeline", "zeitstrahl": "timeline", "roadmap": "timeline", "steps": "timeline",
    "schritte": "timeline", "process": "timeline", "prozess": "timeline", "phasen": "timeline",
    "chart": "chart", "diagramm": "chart", "graph": "chart", "grafik": "chart", "daten": "chart",
    "image": "image", "bild": "image", "foto": "image", "picture": "image",
    "closing": "closing", "abschluss": "closing", "ende": "closing", "thanks": "closing",
    "danke": "closing", "kontakt": "closing", "end": "closing", "outro": "closing",
}

KEY_ALIASES = {
    "heading": "title", "headline": "title", "überschrift": "title", "ueberschrift": "title",
    "name": "title", "kopf": "title",
    "sub": "subtitle", "untertitel": "subtitle", "subhead": "subtitle", "lead": "subtitle",
    "intro": "subtitle", "einleitung": "subtitle", "description": "subtitle",
    "beschreibung": "subtitle", "teaser": "subtitle", "body": "subtitle",
    "points": "bullets", "items": "bullets", "punkte": "bullets", "list": "bullets",
    "liste": "bullets", "content": "bullets", "inhalte": "bullets",
    "eyebrow": "kicker", "label": "kicker", "tag": "kicker", "badge": "kicker",
    "speaker_notes": "notes", "notizen": "notes", "note": "notes",
    "zitat": "quote",
    "typ": "type", "art": "type", "folientyp": "type", "vorlage": "type",
    "author": "attribution", "quelle": "attribution", "source": "attribution",
    "autor": "attribution", "von": "attribution",
    "boxes": "cards", "karten": "cards", "kacheln": "cards", "tiles": "cards",
    "kennzahlen": "stats", "kpis": "stats", "metrics": "stats", "zahlen": "stats",
    "spalten": "columns", "seiten": "columns", "sides": "columns",
    "schritte": "steps", "phasen": "steps", "milestones": "steps", "items_timeline": "steps",
    "spaltennamen": "columns", "header": "columns", "kopfzeile": "columns",
    "zeilen": "rows", "data": "rows", "daten": "rows",
    "kategorien": "categories", "labels": "categories", "x": "categories",
    "reihen": "series", "datenreihen": "series", "values": "series",
    "diagrammtyp": "chart_type", "charttyp": "chart_type", "chart_typ": "chart_type",
    "charttype": "chart_type", "diagrammart": "chart_type",
    "fazit": "takeaway", "insight": "takeaway", "kernaussage": "takeaway", "message": "takeaway",
    "bild": "image_url", "url": "image_url", "img": "image_url", "bild_url": "image_url",
    "kontakt": "contact", "footer": "contact",
}


def _strip_fences(text: str) -> str:
    t = (text or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", t)
        t = re.sub(r"\s*```$", "", t)
    return t.strip()


def parse_slides(raw: Any) -> List[Dict[str, Any]]:
    """Akzeptiert Liste, JSON-String, JSON in Code-Fences oder {'slides': [...]}."""
    data = raw
    if isinstance(raw, str):
        text = _strip_fences(raw)
        if not text:
            return []
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            m = re.search(r"(\[.*\]|\{.*\})", text, re.S)
            if not m:
                raise ValueError(
                    "Der Parameter 'slides' enthält kein gültiges JSON. "
                    "Erwartet wird eine JSON-Liste von Folien-Objekten."
                )
            data = json.loads(m.group(1))
    if isinstance(data, dict):
        for key in ("slides", "folien", "items", "presentation", "praesentation"):
            if isinstance(data.get(key), list):
                data = data[key]
                break
        else:
            data = [data]
    if not isinstance(data, list):
        raise ValueError("'slides' muss eine JSON-Liste von Folien-Objekten sein.")
    return [_normalize_slide(s) for s in data if isinstance(s, (dict, str))]


def _normalize_slide(slide: Any) -> Dict[str, Any]:
    if isinstance(slide, str):
        return {"type": "bullets", "title": slide, "bullets": []}
    out: Dict[str, Any] = {}
    for k, v in slide.items():
        key = str(k).strip().lower().replace(" ", "_").replace("-", "_")
        key = KEY_ALIASES.get(key, key)
        if key in out and not out[key]:
            out[key] = v
        elif key not in out:
            out[key] = v
    raw_type = str(out.get("type") or out.get("layout") or "").strip().lower()
    raw_type = re.sub(r"[\s\-]+", "_", raw_type)
    out["type"] = TYPE_ALIASES.get(raw_type, raw_type if raw_type in TYPE_ALIASES.values() else "")
    if not out["type"]:
        # Typ aus dem Inhalt erraten
        if out.get("quote"):
            out["type"] = "quote"
        elif out.get("stats"):
            out["type"] = "stats"
        elif out.get("cards"):
            out["type"] = "cards"
        elif out.get("steps"):
            out["type"] = "timeline"
        elif out.get("rows"):
            out["type"] = "table"
        elif out.get("series"):
            out["type"] = "chart"
        elif out.get("columns") and isinstance(out.get("columns"), list) and out["columns"] and isinstance(out["columns"][0], dict):
            out["type"] = "compare"
        elif out.get("image_url"):
            out["type"] = "image"
        else:
            out["type"] = "bullets"
    return out


def _as_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, (list, tuple)):
        return " ".join(_as_text(v) for v in value)
    return str(value).strip()


def _bullet_items(value: Any) -> List[Dict[str, Any]]:
    """Normalisiert Bullets zu [{text, level, title}]."""
    items: List[Dict[str, Any]] = []
    if value is None:
        return items
    if isinstance(value, str):
        value = [line for line in value.split("\n") if line.strip()]
    if isinstance(value, dict):
        value = [value]
    for entry in value or []:
        if isinstance(entry, dict):
            text = _as_text(entry.get("text") or entry.get("title") or entry.get("label"))
            head = _as_text(entry.get("title") if entry.get("text") else "")
            level = int(entry.get("level", 0) or 0)
        else:
            text = _as_text(entry)
            head, level = "", 0
            m = re.match(r"^(\s{2,}|\t|[-*•]\s+[-*•]\s*)", text)
            if m:
                level = 1
            text = re.sub(r"^[\s\t]*[-*•]\s*", "", text).strip()
        if not text and not head:
            continue
        items.append({"text": text, "title": head, "level": max(0, min(1, level))})
    return items


# =============================================================================
#  Der Renderer
# =============================================================================


class DeckDesigner:
    """Baut aus einer Folien-Spezifikation eine gestaltete PPTX-Datei."""

    def __init__(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        theme: str = "midnight",
        subtitle: str = "",
        author: str = "",
        date_text: str = "",
        mode: str = "auto",
        font_heading: str = "Cambria",
        font_body: str = "Calibri",
        footer_text: str = "",
        show_footer: bool = True,
        show_slide_numbers: bool = True,
        brand: Optional[Dict[str, str]] = None,
        logo_bytes: Optional[bytes] = None,
    ):
        self.deck_title = _as_text(title) or "Präsentation"
        self.deck_subtitle = _as_text(subtitle)
        self.author = _as_text(author)
        self.date_text = _as_text(date_text)
        self.mode = (mode or "auto").lower()
        self.fh = font_heading or "Cambria"
        self.fb = font_body or "Calibri"
        self.footer_text = _as_text(footer_text)
        self.show_footer = show_footer
        self.show_numbers = show_slide_numbers
        self.logo_bytes = logo_bytes
        self.warnings: List[str] = []

        self.theme_key = resolve_theme(theme)
        self.t = dict(THEMES[self.theme_key])
        self._apply_brand(brand or {})

        self.slides = slides
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.blank = self.prs.slide_layouts[6]
        self._number = 0
        self._section_index = 0

    # ---------------------------------------------------------------- Branding
    def _apply_brand(self, brand: Dict[str, str]) -> None:
        acc = _clean_hex(brand.get("accent", ""))
        acc2 = _clean_hex(brand.get("accent2", ""))
        deep = _clean_hex(brand.get("deep", ""))
        if acc:
            self.t["accent"] = acc
            self.t["series"] = [acc] + [c for c in self.t["series"] if c != acc][:5]
        if acc2:
            self.t["accent2"] = acc2
            self.t["series"] = [self.t["series"][0], acc2] + self.t["series"][1:][:4]
        if deep:
            self.t["deep"] = deep
            self.t["deep2"] = _mix(deep, "FFFFFF", 0.10)
            self.t["on_deep"] = _readable_on(deep)
            self.t["on_deep_mut"] = _mix(self.t["on_deep"], deep, 0.32)
        # Akzent auf beiden Hintergründen lesbar halten
        self.t["accent_on_light"] = _ensure_contrast(self.t["accent"], self.t["light"], 3.2)
        self.t["accent_on_deep"] = _ensure_contrast(self.t["accent"], self.t["deep"], 3.2)
        self.t["accent2_on_light"] = _ensure_contrast(self.t["accent2"], self.t["light"], 3.2)

    # ------------------------------------------------------------- Basis-Utils
    def _new_slide(self, dark: bool):
        s = self.prs.slides.add_slide(self.blank)
        bg = self.t["deep"] if dark else self.t["light"]
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _rgb(bg)
        return s

    def _is_dark(self, kind: str) -> bool:
        if self.mode == "dark":
            return True
        if self.mode == "light":
            return False
        return kind in ("title", "section", "closing", "quote")

    def _pal(self, dark: bool) -> Dict[str, str]:
        t = self.t
        if dark:
            return {
                "bg": t["deep"], "surface": t["deep2"], "surface2": _mix(t["deep"], "FFFFFF", 0.14),
                "ink": t["on_deep"], "muted": t["on_deep_mut"],
                "hair": _mix(t["deep"], t["on_deep"], 0.22), "accent": t["accent_on_deep"],
                "accent2": _ensure_contrast(t["accent2"], t["deep"], 3.2),
            }
        return {
            "bg": t["light"], "surface": t["surface"], "surface2": t["surface2"],
            "ink": t["ink"], "muted": t["muted"], "hair": t["hair"],
            "accent": t["accent_on_light"], "accent2": t["accent2_on_light"],
        }

    def _text(
        self, slide, left, top, width, height, text, *,
        size=16, bold=False, color="000000", font=None, align=PP_ALIGN.LEFT,
        line=1.16, space_after=0, anchor=MSO_ANCHOR.TOP, italic=False,
        spacing=None, caps=False, shrink=True, min_size=9, wrap=True,
    ):
        """Setzt einen Textblock mit optionalem Auto-Shrink gegen Überlauf."""
        text = _as_text(text)
        if not text:
            return None
        if caps:
            text = text.upper()
        if shrink:
            size = self._fit(text, width, height, size, min_size, bold)
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        for i, chunk in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line
            if space_after:
                p.space_after = Pt(space_after)
            r = p.add_run()
            r.text = chunk
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.name = font or self.fb
            f.color.rgb = _rgb(color)
            if spacing:
                try:
                    f._rPr.set("spc", str(int(spacing * 100)))
                except Exception:
                    pass
        return box

    # -- Textmaße: Zeichenbreiten-Schätzung + Umbruch-Simulation ---------------
    @staticmethod
    def _char_em(ch: str, bold: bool) -> float:
        if ch in "iljI|!.,;:'`ı":
            w = 0.30
        elif ch in "ftr()[]{}-/\\":
            w = 0.39
        elif ch in "mwMW@%&€$":
            w = 0.92
        elif ch in "+=<>~^":
            w = 0.60
        elif ch == " ":
            w = 0.28
        elif ch.isdigit():
            w = 0.56
        elif ch.isupper():
            w = 0.67
        else:
            w = 0.53
        return w * (1.06 if bold else 1.0)

    @classmethod
    def _wrap_lines(cls, text: str, width: float, size: float, bold: bool = False) -> int:
        """Simuliert den Zeilenumbruch und liefert die Zeilenzahl."""
        if not text or width <= 0 or size <= 0:
            return 1
        max_em = width * 72.0 / size
        total = 0
        for para in str(text).split("\n"):
            cur, count = 0.0, 0
            for word in para.split(" "):
                w = sum(cls._char_em(c, bold) for c in word)
                add = w if count == 0 else w + cls._char_em(" ", bold)
                if count and cur + add > max_em:
                    total += 1
                    cur, count = w, 1
                else:
                    cur += add
                    count += 1
            total += 1
        return max(1, total)

    def _measure(self, text: str, width: float, base: float, bold: bool = False,
                 line: float = 1.18, max_lines: Optional[int] = None,
                 min_size: float = 9) -> Tuple[float, int, float]:
        """Ermittelt (Schriftgröße, Zeilen, Höhe) unter einer Zeilenbegrenzung."""
        size = float(base)
        while size > min_size:
            n = self._wrap_lines(text, width, size, bold)
            if max_lines is None or n <= max_lines:
                return size, n, n * size * line / 72.0
            size -= 1
        n = self._wrap_lines(text, width, min_size, bold)
        return min_size, n, n * min_size * line / 72.0

    def _fit(self, text: str, width: float, height: float, base: float,
             min_size: float = 9, bold: bool = False) -> float:
        """Größte Schriftgröße, mit der der Text in den Rahmen passt."""
        size = float(base)
        while size > min_size:
            lines = self._wrap_lines(text, width, size, bold)
            if lines * (size * 1.24 / 72.0) <= height + 0.03:
                return size
            size -= 1
        return min_size

    def _shape(self, slide, kind, left, top, width, height, fill=None,
               line_color=None, line_w=0.75, radius=None, shadow=False):
        sh = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
        # Theme-Stil (harter Standardschatten, Themefarben) entfernen
        try:
            from pptx.oxml.ns import qn

            style = sh._element.find(qn("p:style"))
            if style is not None:
                sh._element.remove(style)
        except Exception:
            pass
        if radius is not None:
            try:
                sh.adjustments[0] = radius
            except Exception:
                pass
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = _rgb(fill)
        else:
            sh.fill.background()
        if line_color:
            sh.line.color.rgb = _rgb(line_color)
            sh.line.width = Pt(line_w)
        else:
            sh.line.fill.background()
        try:
            sh.shadow.inherit = False  # leeres <a:effectLst/> anlegen
        except Exception:
            pass
        if shadow:
            self._soft_shadow(sh)
        sh.text_frame.word_wrap = True
        return sh

    def _soft_shadow(self, shape, blur_pt: float = 16, dist_pt: float = 3.5,
                     alpha: float = 0.11):
        """Sehr weicher, dezenter Schlagschatten (statt PowerPoint-Standard)."""
        try:
            from lxml import etree
            from pptx.oxml.ns import qn

            spPr = shape._element.spPr
            eff = spPr.find(qn("a:effectLst"))
            if eff is None:
                eff = etree.SubElement(spPr, qn("a:effectLst"))
            for child in list(eff):
                eff.remove(child)
            shdw = etree.SubElement(eff, qn("a:outerShdw"))
            shdw.set("blurRad", str(int(blur_pt * 12700)))
            shdw.set("dist", str(int(dist_pt * 12700)))
            shdw.set("dir", "5400000")
            shdw.set("rotWithShape", "0")
            clr = etree.SubElement(shdw, qn("a:srgbClr"))
            clr.set("val", "000000")
            al = etree.SubElement(clr, qn("a:alpha"))
            al.set("val", str(int(max(0.0, min(1.0, alpha)) * 100000)))
        except Exception:
            pass

    def _card(self, slide, left, top, width, height, fill, border=None, shadow=True):
        return self._shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                           fill=fill, line_color=border, radius=0.05, shadow=shadow)

    def _circle_number(self, slide, left, top, size, label, fill, text_color, font_size=None):
        c = self._shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill=fill)
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(label)
        r.font.size = Pt(font_size or max(9, size * 26))
        r.font.bold = True
        r.font.name = self.fb
        r.font.color.rgb = _rgb(text_color)
        return c

    # -------------------------------------------------------------- Dekoration
    def _stage_decor(self, slide, variant: int = 0):
        """Wiederkehrendes Motiv: weiche konzentrische Kreise + Punktraster."""
        t = self.t
        base, on = t["deep"], t["on_deep"]
        glow = _mix(base, t["accent"], 0.30)
        soft = _mix(base, on, 0.07)
        if variant % 2 == 0:
            self._shape(slide, MSO_SHAPE.OVAL, 9.05, -1.65, 6.2, 6.2, fill=_mix(base, glow, 0.55))
            self._shape(slide, MSO_SHAPE.OVAL, 10.15, -0.55, 4.0, 4.0, fill=soft)
            self._shape(slide, MSO_SHAPE.OVAL, 8.15, 1.55, 3.1, 3.1, fill=None,
                        line_color=_mix(base, on, 0.30), line_w=1.0)
            self._dot_grid(slide, 10.55, 5.15, 6, 3, 0.30, _mix(base, on, 0.30))
        else:
            self._shape(slide, MSO_SHAPE.OVAL, -1.9, 3.4, 5.6, 5.6, fill=_mix(base, glow, 0.45))
            self._shape(slide, MSO_SHAPE.OVAL, 10.9, -1.4, 4.4, 4.4, fill=soft)
            self._dot_grid(slide, 11.15, 5.55, 5, 3, 0.30, _mix(base, on, 0.28))

    def _dot_grid(self, slide, left, top, cols, rows, step, color, dot=0.062):
        for r in range(rows):
            for c in range(cols):
                self._shape(slide, MSO_SHAPE.OVAL, left + c * step, top + r * step,
                            dot, dot, fill=color)

    def _footer(self, slide, dark: bool, number: Optional[int]):
        if not self.show_footer and not self.show_numbers:
            return
        pal = self._pal(dark)
        col = _mix(pal["muted"], pal["bg"], 0.25)
        left_txt = self.footer_text or self.deck_title
        if self.show_footer and left_txt:
            self._text(slide, MX, FOOTER_Y, CW - 1.0, 0.28, left_txt[:78],
                       size=9, color=col, shrink=False)
        if self.show_numbers and number:
            self._text(slide, SW - MX - 1.0, FOOTER_Y, 1.0, 0.28, str(number),
                       size=9, color=col, align=PP_ALIGN.RIGHT, shrink=False)

    def _slide_head(self, slide, dark, title, subtitle="", kicker=""):
        """Kopfbereich für Inhaltsfolien. Gibt das obere Y des Inhaltsbereichs zurück."""
        pal = self._pal(dark)
        y = TITLE_Y
        if kicker:
            self._text(slide, MX, y, CW, 0.24, kicker, size=10.5, bold=True,
                       color=pal["accent"], caps=True, spacing=1.4, shrink=False)
            y += 0.36
        title = _as_text(title)
        if title:
            tw = CW - 0.4
            size, lines, h = self._measure(title, tw, 31, bold=True, line=1.14,
                                           max_lines=2, min_size=20)
            self._text(slide, MX, y, tw, h + 0.10, title, size=size, bold=True,
                       color=pal["ink"], font=self.fh, line=1.14, shrink=False)
            y += h + 0.20
        sub = _as_text(subtitle)
        if sub:
            sw = CW - 1.6
            size, lines, h = self._measure(sub, sw, 15, line=1.24, max_lines=2, min_size=11.5)
            self._text(slide, MX, y, sw, h + 0.08, sub, size=size,
                       color=pal["muted"], line=1.24, shrink=False)
            y += h + 0.34
        else:
            y += 0.14
        return y

    # =========================================================== Folien-Layouts
    def _slide_title(self, s: Dict[str, Any]):
        dark = self._is_dark("title")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        if dark:
            self._stage_decor(slide, 0)
        else:
            self._shape(slide, MSO_SHAPE.OVAL, 9.4, -1.5, 5.8, 5.8,
                        fill=_mix(self.t["light"], self.t["accent"], 0.10))
            self._shape(slide, MSO_SHAPE.OVAL, 8.5, 2.0, 3.0, 3.0, fill=None,
                        line_color=_mix(self.t["light"], self.t["accent"], 0.28), line_w=1.0)
            self._dot_grid(slide, 10.6, 5.2, 6, 3, 0.30,
                           _mix(self.t["light"], self.t["accent"], 0.35))

        kicker = _as_text(s.get("kicker")) or _as_text(s.get("category"))
        title = _as_text(s.get("title")) or self.deck_title
        sub = _as_text(s.get("subtitle")) or self.deck_subtitle

        # Blockhöhe messen und den Satz vertikal ausbalancieren
        tw, sw = 8.30, 7.40
        t_size, t_lines, t_h = self._measure(title, tw, 46, bold=True, line=1.10,
                                             max_lines=3, min_size=26)
        s_size = s_h = 0.0
        if sub:
            s_size, _s_lines, s_h = self._measure(sub, sw, 17, line=1.30, max_lines=3,
                                                  min_size=13)
        block = (0.46 if kicker else 0.0) + t_h + (s_h + 0.30 if sub else 0.0)
        y = max(1.75, min(2.55, (SH - 1.0 - block) / 2))

        if kicker:
            self._circle_number(slide, MX, y + 0.03, 0.19, "", pal["accent"], pal["bg"])
            self._text(slide, MX + 0.34, y, 7.0, 0.26, kicker, size=11.5, bold=True,
                       color=pal["accent"], caps=True, spacing=1.6, shrink=False)
            y += 0.46
        self._text(slide, MX, y, tw, t_h + 0.12, title, size=t_size, bold=True,
                   color=pal["ink"], font=self.fh, line=1.10, shrink=False)
        y += t_h + 0.30 + t_size / 190.0
        if sub:
            self._text(slide, MX, y, sw, s_h + 0.10, sub, size=s_size, color=pal["muted"],
                       line=1.30, shrink=False)

        meta = " · ".join([x for x in [self.author, self.date_text] if x])
        if meta:
            self._text(slide, MX, 6.62, 8.0, 0.32, meta, size=11.5,
                       color=_mix(pal["muted"], pal["bg"], 0.15), shrink=False)
        self._place_logo(slide, dark)
        return slide

    def _slide_section(self, s: Dict[str, Any]):
        self._section_index += 1
        dark = self._is_dark("section")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        if dark:
            self._stage_decor(slide, 1)
        else:
            self._shape(slide, MSO_SHAPE.OVAL, -1.6, 3.6, 5.2, 5.2,
                        fill=_mix(self.t["light"], self.t["accent"], 0.09))

        idx = s.get("index") or s.get("number") or self._section_index
        try:
            idx_txt = f"{int(idx):02d}"
        except (TypeError, ValueError):
            idx_txt = _as_text(idx)[:2]

        title = _as_text(s.get("title"))
        sub = _as_text(s.get("subtitle"))
        t_size, _n, t_h = self._measure(title, 9.0, 40, bold=True, line=1.08,
                                        max_lines=2, min_size=26)
        s_h = 0.0
        if sub:
            s_size, _n2, s_h = self._measure(sub, 7.4, 15.5, line=1.28, max_lines=2,
                                             min_size=12)
        num_h = 1.24
        block = num_h + 0.30 + t_h + (s_h + 0.28 if sub else 0.0)
        y = (SH - 0.8 - block) / 2

        self._text(slide, MX, y - 0.14, 3.2, num_h + 0.3, idx_txt, size=78, bold=True,
                   color=_mix(pal["bg"], pal["accent"], 0.60), font=self.fh, shrink=False)
        y += num_h + 0.30
        self._text(slide, MX, y, 9.0, t_h + 0.12, title, size=t_size, bold=True,
                   color=pal["ink"], font=self.fh, line=1.08, shrink=False)
        y += t_h + 0.28
        if sub:
            self._text(slide, MX, y, 7.4, s_h + 0.10, sub, size=s_size, color=pal["muted"],
                       line=1.28, shrink=False)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_bullets(self, s: Dict[str, Any]):
        dark = self._is_dark("bullets")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))
        items = _bullet_items(s.get("bullets"))
        if not items:
            self._footer(slide, dark, self._number)
            return slide

        avail = BODY_BOTTOM - y
        tops = [i for i in items if i["level"] == 0]
        # Wenige Hauptpunkte -> gestaltete Zeilen-Karten, sonst kompakte Liste
        if len(items) <= 5 and len(tops) == len(items) and avail > 3.2:
            gap = 0.22
            h = min(1.16, (avail - gap * (len(items) - 1)) / len(items))
            y += min(0.50, max(0.0, (avail - (h * len(items) + gap * (len(items) - 1))) / 2))
            for i, it in enumerate(items):
                top = y + i * (h + gap)
                self._card(slide, MX, top, CW, h,
                           pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.08))
                self._circle_number(slide, MX + 0.42, top + (h - 0.46) / 2, 0.46, i + 1,
                                    pal["accent"], _readable_on(pal["accent"]), font_size=14)
                tx, tw = MX + 1.18, CW - 1.18 - 0.55
                head, body = it["title"], it["text"]
                if head and body:
                    hs, _n, hh = self._measure(head, tw, 16.5, bold=True, line=1.14,
                                               max_lines=1, min_size=12)
                    bs, _n2, bh = self._measure(body, tw, 13.5, line=1.24, max_lines=2,
                                                min_size=10)
                    gy = top + max(0.10, (h - (hh + 0.10 + bh)) / 2)
                    self._text(slide, tx, gy, tw, hh + 0.06, head, size=hs, bold=True,
                               color=pal["ink"], font=self.fh, line=1.14, shrink=False)
                    self._text(slide, tx, gy + hh + 0.10, tw, bh + 0.06, body, size=bs,
                               color=pal["muted"], line=1.24, shrink=False)
                else:
                    self._text(slide, tx, top, tw, h, head or body, size=16,
                               color=pal["ink"], line=1.24, anchor=MSO_ANCHOR.MIDDLE,
                               min_size=11)
        else:
            n = len(items)
            base = 17.0 if n <= 6 else (15.0 if n <= 9 else 13.0)
            rows, total = [], 0.0
            for _attempt in range(12):
                rows, total = [], 0.0
                for i, it in enumerate(items):
                    lvl = it["level"]
                    size = base if lvl == 0 else max(10.0, base - 2.5)
                    indent = 0.0 if lvl == 0 else 0.58
                    tw = CW - indent - 0.45
                    txt = (f"{it['title']} — {it['text']}" if it["title"] and it["text"]
                           else (it["title"] or it["text"]))
                    h = self._wrap_lines(txt, tw, size) * size * 1.26 / 72.0
                    gap = 0.0 if i == 0 else (0.22 if lvl == 0 else 0.09)
                    rows.append((txt, size, indent, tw, h, gap, lvl, it))
                    total += h + gap
                if total <= avail or base <= 10.5:
                    break
                base -= 1
            lead = min(0.26, max(0.0, (avail - total) / max(1, len(rows) - 1)))
            top = y
            for txt, size, indent, tw, h, gap, lvl, it in rows:
                top += gap + (lead if gap else 0.0)
                if top + h > BODY_BOTTOM + 0.12:
                    break
                dot = 0.115 if lvl == 0 else 0.078
                self._shape(slide, MSO_SHAPE.OVAL, MX + indent + 0.03,
                            top + (size * 1.26 / 72.0 - dot) / 2, dot, dot,
                            fill=pal["accent"] if lvl == 0
                            else _mix(pal["accent"], pal["bg"], 0.42))
                self._text(slide, MX + indent + 0.42, top, tw, h + 0.06, txt, size=size,
                           color=pal["ink"] if lvl == 0 else pal["muted"], line=1.26,
                           bold=(lvl == 0 and bool(it["title"]) and not it["text"]),
                           shrink=False)
                top += h
        self._footer(slide, dark, self._number)
        return slide

    def _slide_agenda(self, s: Dict[str, Any]):
        dark = self._is_dark("agenda")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title") or "Agenda", s.get("subtitle"),
                             s.get("kicker"))
        items = _bullet_items(s.get("bullets"))[:10]
        if not items:
            self._footer(slide, dark, self._number)
            return slide
        two_col = len(items) > 5
        per = math.ceil(len(items) / 2) if two_col else len(items)
        col_w = (CW - 0.6) / 2 if two_col else CW
        avail = BODY_BOTTOM - y
        h = min(0.82, avail / max(per, 1))
        y += min(0.30, max(0.0, (avail - h * per) / 2))
        for i, it in enumerate(items):
            c, r = (i // per, i % per) if two_col else (0, i)
            left = MX + c * (col_w + 0.6)
            top = y + r * h
            self._text(slide, left, top, 0.62, h - 0.08, f"{i + 1:02d}", size=21, bold=True,
                       color=_mix(pal["accent"], pal["bg"], 0.32), font=self.fh,
                       anchor=MSO_ANCHOR.MIDDLE, shrink=False)
            txt = f"{it['title']} — {it['text']}" if it["title"] and it["text"] else (it["title"] or it["text"])
            self._text(slide, left + 0.72, top, col_w - 0.78, h - 0.08, txt, size=16.5,
                       color=pal["ink"], anchor=MSO_ANCHOR.MIDDLE, line=1.18, min_size=11)
            if r < per - 1 and i < len(items) - 1:
                self._shape(slide, MSO_SHAPE.RECTANGLE, left, top + h - 0.03,
                            col_w - 0.1, 0.010, fill=pal["hair"])
        self._footer(slide, dark, self._number)
        return slide

    def _slide_cards(self, s: Dict[str, Any]):
        dark = self._is_dark("cards")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))
        raw = s.get("cards") or s.get("bullets") or []
        cards = []
        for c in raw if isinstance(raw, list) else []:
            if isinstance(c, dict):
                cards.append({
                    "title": _as_text(c.get("title") or c.get("label") or c.get("heading")),
                    "text": _as_text(c.get("text") or c.get("description") or c.get("body")),
                    "icon": _as_text(c.get("icon") or c.get("badge")),
                })
            else:
                cards.append({"title": "", "text": _as_text(c), "icon": ""})
        cards = [c for c in cards if c["title"] or c["text"]][:6]
        if not cards:
            self._footer(slide, dark, self._number)
            return slide

        n = len(cards)
        cols = 3 if n in (3, 5, 6) else 2
        if n == 1:
            cols = 1
        if n == 4:
            cols = 2
        rows = math.ceil(n / cols)
        gap = 0.32
        cw = (CW - gap * (cols - 1)) / cols
        avail = BODY_BOTTOM - y
        inner = cw - 0.76
        max_h = (avail - gap * (rows - 1)) / rows

        # Einheitliche Typografie über alle Karten + Kartenhöhe aus dem Inhalt
        t_size = min([self._measure(c["title"], inner, 17, bold=True, line=1.12,
                                    max_lines=2, min_size=12)[0]
                      for c in cards if c["title"]] or [17])
        t_h = max([self._wrap_lines(c["title"], inner, t_size, True) * t_size * 1.12 / 72
                   for c in cards if c["title"]] or [0.0])
        b_size = min([self._measure(c["text"], inner, 13.5, line=1.28,
                                    max_lines=max(2, int((max_h - 1.9) / 0.24) + 3),
                                    min_size=9.5)[0]
                      for c in cards if c["text"]] or [13.5])
        b_h = max([self._wrap_lines(c["text"], inner, b_size, False) * b_size * 1.28 / 72
                   for c in cards if c["text"]] or [0.0])
        need = 1.04 + (t_h + 0.16 if t_h else 0.0) + b_h + 0.38
        ch = max(1.85, min(need, max_h))
        top0 = y + min(0.55, max(0.0, (avail - (ch * rows + gap * (rows - 1))) / 2))

        for i, c in enumerate(cards):
            r, col = divmod(i, cols)
            left = MX + col * (cw + gap)
            top = top0 + r * (ch + gap)
            fill = pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.08)
            self._card(slide, left, top, cw, ch, fill)
            badge = c["icon"] or str(i + 1)
            self._circle_number(slide, left + 0.38, top + 0.36, 0.50, badge[:2],
                                pal["accent"], _readable_on(pal["accent"]), font_size=14)
            ty = top + 1.04
            if c["title"]:
                self._text(slide, left + 0.38, ty, inner, t_h + 0.06, c["title"],
                           size=t_size, bold=True, color=pal["ink"], font=self.fh,
                           line=1.12, shrink=False)
            ty += (t_h + 0.16) if t_h else 0.0
            if c["text"]:
                self._text(slide, left + 0.38, ty, inner, max(0.3, top + ch - 0.30 - ty),
                           c["text"], size=b_size, color=pal["muted"], line=1.28,
                           min_size=9.5)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_stats(self, s: Dict[str, Any]):
        dark = self._is_dark("stats")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))
        raw = s.get("stats") or []
        stats = []
        for st in raw if isinstance(raw, list) else []:
            if isinstance(st, dict):
                stats.append({
                    "value": _as_text(st.get("value") or st.get("wert") or st.get("number")),
                    "label": _as_text(st.get("label") or st.get("title") or st.get("name")),
                    "note": _as_text(st.get("note") or st.get("description") or st.get("text")),
                })
            else:
                stats.append({"value": _as_text(st), "label": "", "note": ""})
        stats = [s2 for s2 in stats if s2["value"] or s2["label"]][:4]
        if not stats:
            self._footer(slide, dark, self._number)
            return slide

        n = len(stats)
        gap = 0.32
        cw = (CW - gap * (n - 1)) / n
        avail = BODY_BOTTOM - y
        fill = pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.08)
        col = _ensure_contrast(pal["accent"], fill, 3.0)
        # Einheitliche Schriftgrößen über alle Kacheln, damit die Zeilen fluchten
        tw_in = cw - 0.68
        vsize = min([self._measure(st["value"], tw_in, 52, bold=True, max_lines=1,
                                   min_size=22)[0] for st in stats])
        lsize = min([self._measure(st["label"], tw_in, 15, bold=True, line=1.16,
                                   max_lines=2, min_size=11)[0]
                     for st in stats if st["label"]] or [15])
        label_h = max([self._wrap_lines(st["label"], tw_in, lsize, True) * lsize * 1.16 / 72
                       for st in stats if st["label"]] or [0.0])
        note_h = max([self._wrap_lines(st["note"], tw_in, 12, False) * 12 * 1.26 / 72
                      for st in stats if st["note"]] or [0.0])
        ch = max(1.75, min(1.44 + label_h + 0.20 + note_h + 0.34, avail))
        top = y + min(0.45, max(0.0, (avail - ch) / 2))
        for i, st in enumerate(stats):
            left = MX + i * (cw + gap)
            self._card(slide, left, top, cw, ch, fill)
            self._text(slide, left + 0.34, top + 0.42, tw_in, 0.92, st["value"],
                       size=vsize, bold=True, color=col, font=self.fh, line=1.0,
                       shrink=False)
            ly = top + 1.44
            if st["label"]:
                self._text(slide, left + 0.34, ly, tw_in, label_h + 0.06, st["label"],
                           size=lsize, bold=True, color=pal["ink"], line=1.16, shrink=False)
            ly += label_h + 0.20
            if st["note"]:
                self._text(slide, left + 0.34, ly, tw_in, max(0.3, top + ch - 0.30 - ly),
                           st["note"], size=12, color=pal["muted"], line=1.26, min_size=9.5)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_quote(self, s: Dict[str, Any]):
        dark = self._is_dark("quote")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        if dark:
            self._shape(slide, MSO_SHAPE.OVAL, 9.6, -1.3, 5.4, 5.4,
                        fill=_mix(pal["bg"], self.t["accent"], 0.22))
            self._shape(slide, MSO_SHAPE.OVAL, -1.5, 4.2, 4.2, 4.2,
                        fill=_mix(pal["bg"], pal["ink"], 0.05))
        else:
            self._shape(slide, MSO_SHAPE.OVAL, 9.9, -1.2, 5.0, 5.0,
                        fill=_mix(pal["bg"], self.t["accent"], 0.08))

        quote = _as_text(s.get("quote") or s.get("text") or s.get("title"))
        attribution = _as_text(s.get("attribution") or s.get("subtitle"))
        qw = 10.5
        q_size, _n, q_h = self._measure(quote, qw, 30, line=1.28, max_lines=5, min_size=16)
        group = 0.85 + q_h + (0.92 if attribution else 0.0)
        top = max(0.75, (SH - 0.55 - group) / 2)

        self._text(slide, MX, top - 0.70, 2.0, 1.4, "“", size=110, bold=True,
                   color=_mix(pal["bg"], pal["accent"], 0.55), font=self.fh, shrink=False)
        q_top = top + 0.85
        self._text(slide, MX + 0.06, q_top, qw, q_h + 0.12, quote, size=q_size,
                   color=pal["ink"], font=self.fh, italic=True, line=1.28, shrink=False)
        if attribution:
            a_top = q_top + q_h * 1.12 + 0.42
            self._circle_number(slide, MX + 0.07, a_top + 0.07, 0.15, "", pal["accent"],
                                pal["bg"])
            self._text(slide, MX + 0.40, a_top, 8.6, 0.42, attribution, size=13.5, bold=True,
                       color=pal["muted"], min_size=10)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_compare(self, s: Dict[str, Any]):
        dark = self._is_dark("compare")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))

        cols = s.get("columns")
        if not isinstance(cols, list) or not cols:
            cols = [c for c in (s.get("left"), s.get("right")) if isinstance(c, dict)]
        norm = []
        for c in (cols or [])[:2]:
            if not isinstance(c, dict):
                continue
            norm.append({
                "title": _as_text(c.get("title") or c.get("heading") or c.get("label")),
                "subtitle": _as_text(c.get("subtitle") or c.get("note")),
                "bullets": _bullet_items(c.get("bullets") or c.get("items") or c.get("points")),
                "highlight": bool(c.get("highlight")),
            })
        if len(norm) < 2:
            return self._slide_bullets(s)
        if not any(c["highlight"] for c in norm):
            norm[1]["highlight"] = True

        gap = 0.42
        cw = (CW - gap) / 2
        avail = BODY_BOTTOM - y
        need = 0.0
        for c in norm:
            hgt = 0.44 + (0.54 if c["title"] else 0.0) + (0.48 if c["subtitle"] else 0.0) + 0.06
            for it in c["bullets"][:7]:
                txt = f"{it['title']} — {it['text']}" if it["title"] and it["text"] else (it["title"] or it["text"])
                hgt += self._wrap_lines(txt, cw - 1.22, 14) * 14 * 1.26 / 72 + 0.26
            need = max(need, hgt + 0.24)
        ch = max(2.2, min(avail, need))
        y += min(0.40, max(0.0, (avail - ch) / 2))
        for i, c in enumerate(norm):
            left = MX + i * (cw + gap)
            if c["highlight"]:
                fill = _mix(pal["bg"], pal["accent"], 0.12 if not dark else 0.22)
            else:
                fill = pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.07)
            self._card(slide, left, y, cw, ch, fill)
            ty = y + 0.44
            head_col = pal["ink"]
            if c["title"]:
                self._text(slide, left + 0.42, ty, cw - 0.84, 0.46, c["title"], size=20,
                           bold=True, color=head_col, font=self.fh, min_size=14)
                ty += 0.54
            if c["subtitle"]:
                self._text(slide, left + 0.42, ty, cw - 0.84, 0.40, c["subtitle"], size=12.5,
                           color=pal["muted"], min_size=10)
                ty += 0.48
            ty += 0.06
            items = c["bullets"][:7]
            for it in items:
                txt = f"{it['title']} — {it['text']}" if it["title"] and it["text"] else (it["title"] or it["text"])
                tw = cw - 0.84 - 0.38
                size = self._fit(txt, tw, 1.2, 14, 10)
                lines = max(1, math.ceil(len(txt) / max(1, int(tw / (0.5 * size / 72)))))
                h = lines * (size * 1.26 / 72) + 0.06
                if ty + h > y + ch - 0.30:
                    break
                self._shape(slide, MSO_SHAPE.OVAL, left + 0.44, ty + size / 72 * 0.40,
                            0.10, 0.10, fill=pal["accent"])
                self._text(slide, left + 0.78, ty, tw, h, txt, size=size, color=pal["ink"],
                           line=1.26, shrink=False)
                ty += h + 0.20
        self._footer(slide, dark, self._number)
        return slide

    def _slide_table(self, s: Dict[str, Any]):
        dark = self._is_dark("table")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))

        cols = s.get("columns") or []
        if cols and isinstance(cols[0], dict):
            cols = [_as_text(c.get("title") or c.get("label")) for c in cols]
        cols = [_as_text(c) for c in cols][:7]
        rows = s.get("rows") or []
        clean_rows = []
        for r in rows if isinstance(rows, list) else []:
            if isinstance(r, dict):
                clean_rows.append([_as_text(r.get(c, "")) for c in cols])
            elif isinstance(r, (list, tuple)):
                clean_rows.append([_as_text(v) for v in r][:len(cols) or 7])
        clean_rows = clean_rows[:9]
        if not cols and clean_rows:
            cols = [""] * len(clean_rows[0])
        if not cols:
            self._footer(slide, dark, self._number)
            return slide

        n_rows = len(clean_rows) + 1
        avail = BODY_BOTTOM - y
        head_h = 0.50
        row_h = min(0.60, max(0.34, (avail - head_h) / max(len(clean_rows), 1)))
        total_h = head_h + row_h * len(clean_rows)
        y += min(0.35, max(0.0, (avail - total_h) / 2))
        gt = slide.shapes.add_table(n_rows, len(cols), Inches(MX), Inches(y),
                                    Inches(CW), Inches(total_h))
        table = gt.table
        try:
            table.first_row = False
            table.horz_banding = False
        except Exception:
            pass
        table.rows[0].height = Inches(head_h)
        for r in range(1, n_rows):
            table.rows[r].height = Inches(row_h)

        head_fill = pal["accent"]
        head_ink = _readable_on(head_fill)
        fsize = 13 if len(cols) <= 4 else (11.5 if len(cols) <= 5 else 10.5)

        def style_cell(cell, text, *, bold=False, color="000000", fill=None, size=12,
                       align=PP_ALIGN.LEFT):
            cell.margin_left = Inches(0.16)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(fill)
            else:
                cell.fill.background()
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = _as_text(text)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = self.fb
            r.font.color.rgb = _rgb(color)

        for c, name in enumerate(cols):
            style_cell(table.cell(0, c), name, bold=True, color=head_ink, fill=head_fill,
                       size=fsize)
        zebra = pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.07)
        for r, row in enumerate(clean_rows, start=1):
            for c in range(len(cols)):
                val = row[c] if c < len(row) else ""
                style_cell(table.cell(r, c), val,
                           bold=(c == 0),
                           color=pal["ink"] if c == 0 else pal["muted"],
                           fill=zebra if r % 2 == 1 else pal["bg"],
                           size=fsize,
                           align=PP_ALIGN.LEFT)
        note = _as_text(s.get("note") or s.get("source"))
        if note and y + total_h + 0.32 < FOOTER_Y:
            self._text(slide, MX, y + total_h + 0.16, CW, 0.28, note, size=10,
                       color=pal["muted"], shrink=False)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_timeline(self, s: Dict[str, Any]):
        dark = self._is_dark("timeline")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))
        raw = s.get("steps") or s.get("bullets") or []
        steps = []
        for st in raw if isinstance(raw, list) else []:
            if isinstance(st, dict):
                steps.append({
                    "label": _as_text(st.get("label") or st.get("date") or st.get("phase")),
                    "title": _as_text(st.get("title") or st.get("name") or st.get("heading")),
                    "text": _as_text(st.get("text") or st.get("description")),
                })
            else:
                steps.append({"label": "", "title": _as_text(st), "text": ""})
        steps = [x for x in steps if x["title"] or x["text"] or x["label"]][:5]
        if not steps:
            self._footer(slide, dark, self._number)
            return slide

        n = len(steps)
        avail = BODY_BOTTOM - y
        seg = CW / n
        col_w = seg - 0.24
        dia = 0.66 if avail > 3.6 else 0.58
        has_label = any(st["label"] for st in steps)

        t_size = min([self._measure(st["title"], col_w, 17, bold=True, line=1.12,
                                    max_lines=2, min_size=11)[0]
                      for st in steps if st["title"]] or [16])
        t_h = max([self._wrap_lines(st["title"], col_w, t_size, True) * t_size * 1.12 / 72
                   for st in steps if st["title"]] or [0.0])
        b_size = min([self._measure(st["text"], col_w, 12.5, line=1.28, max_lines=4,
                                    min_size=9.5)[0]
                      for st in steps if st["text"]] or [12.5])
        b_h = max([self._wrap_lines(st["text"], col_w, b_size, False) * b_size * 1.28 / 72
                   for st in steps if st["text"]] or [0.0])

        head = (0.34 if has_label else 0.0) + t_h + 0.20
        group = head + dia + (0.26 + b_h if b_h else 0.0)
        top = y + max(0.0, (avail - group) / 2)
        axis_y = top + head + dia / 2

        centers = [MX + seg * (i + 0.5) for i in range(n)]
        self._shape(slide, MSO_SHAPE.RECTANGLE, centers[0], axis_y - 0.008,
                    centers[-1] - centers[0], 0.016, fill=_mix(pal["hair"], pal["bg"], 0.10))
        for i, st in enumerate(steps):
            cx = centers[i]
            ty = top
            if has_label:
                self._text(slide, cx - col_w / 2, ty, col_w, 0.28, st["label"], size=11.5,
                           bold=True, color=pal["accent"], align=PP_ALIGN.CENTER, caps=True,
                           spacing=1.0, shrink=False)
                ty += 0.34
            if st["title"]:
                self._text(slide, cx - col_w / 2, ty, col_w, t_h + 0.06, st["title"],
                           size=t_size, bold=True, color=pal["ink"], font=self.fh,
                           align=PP_ALIGN.CENTER, line=1.12, shrink=False)
            self._circle_number(slide, cx - dia / 2, axis_y - dia / 2, dia, i + 1,
                                pal["accent"], _readable_on(pal["accent"]), font_size=15)
            if st["text"]:
                self._text(slide, cx - col_w / 2, axis_y + dia / 2 + 0.26, col_w, b_h + 0.08,
                           st["text"], size=b_size, color=pal["muted"], align=PP_ALIGN.CENTER,
                           line=1.28, shrink=False)
        self._footer(slide, dark, self._number)
        return slide

    def _slide_chart(self, s: Dict[str, Any]):
        dark = self._is_dark("chart")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        y = self._slide_head(slide, dark, s.get("title"), s.get("subtitle"), s.get("kicker"))

        cats = s.get("categories") or []
        series_in = s.get("series") or []
        if isinstance(series_in, dict):
            series_in = [{"name": k, "values": v} for k, v in series_in.items()]
        series = []
        for ser in series_in if isinstance(series_in, list) else []:
            if isinstance(ser, dict):
                name = _as_text(ser.get("name") or ser.get("label") or "Serie")
                vals = ser.get("values") or ser.get("data") or []
            else:
                name, vals = "Serie", ser
            clean = []
            for v in vals if isinstance(vals, (list, tuple)) else []:
                try:
                    clean.append(float(str(v).replace("%", "").replace(",", ".").strip()))
                except (TypeError, ValueError):
                    clean.append(0.0)
            if clean:
                series.append((name, clean))
        cats = [_as_text(c) for c in cats] if isinstance(cats, list) else []
        if not series:
            self.warnings.append("Diagramm ohne Datenreihen übersprungen.")
            self._footer(slide, dark, self._number)
            return slide
        if not cats:
            cats = [str(i + 1) for i in range(len(series[0][1]))]

        takeaway = _as_text(s.get("takeaway"))
        avail_h = BODY_BOTTOM - y
        if takeaway:
            tw = 3.65
            self._card(slide, MX, y, tw, avail_h,
                       pal["surface"] if not dark else _mix(pal["bg"], pal["ink"], 0.07))
            k_size, _n, k_h = self._measure(takeaway, tw - 0.68, 15.5, line=1.32,
                                            max_lines=9, min_size=11)
            ty = y + max(0.36, (avail_h - (k_h + 0.46)) / 2)
            self._text(slide, MX + 0.34, ty, tw - 0.68, 0.26, "Kernaussage", size=10.5,
                       bold=True, color=pal["accent"], caps=True, spacing=1.4, shrink=False)
            self._text(slide, MX + 0.34, ty + 0.46, tw - 0.68, k_h + 0.10, takeaway,
                       size=k_size, color=pal["ink"], line=1.32, shrink=False)
            cx, cw = MX + tw + 0.34, CW - tw - 0.34
        else:
            cx, cw = MX - 0.1, CW + 0.2

        kind = str(s.get("chart_type") or "column").strip().lower()
        if kind in ("pie", "kreis", "torte", "doughnut", "donut", "ring") and not takeaway:
            cx, cw = MX + CW * 0.16, CW * 0.68
        kmap = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED, "balken": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED, "säulen": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "saeulen": XL_CHART_TYPE.COLUMN_CLUSTERED, "spalten": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "stacked": XL_CHART_TYPE.COLUMN_STACKED, "gestapelt": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE, "linie": XL_CHART_TYPE.LINE,
            "linien": XL_CHART_TYPE.LINE, "verlauf": XL_CHART_TYPE.LINE,
            "area": XL_CHART_TYPE.AREA, "fläche": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE, "kreis": XL_CHART_TYPE.PIE, "torte": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT, "donut": XL_CHART_TYPE.DOUGHNUT,
            "ring": XL_CHART_TYPE.DOUGHNUT,
        }
        ctype = kmap.get(kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
        is_pie = ctype in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)
        if is_pie:
            series = series[:1]

        cd = CategoryChartData()
        cd.categories = cats
        for name, vals in series:
            cd.add_series(name, [v for v in vals[:len(cats)]] + [0.0] * max(0, len(cats) - len(vals)))

        gf = slide.shapes.add_chart(ctype, Inches(cx), Inches(y - 0.05), Inches(cw),
                                    Inches(avail_h + 0.05), cd)
        chart = gf.chart
        colors = [_ensure_contrast(c, pal["bg"], 1.6) for c in self.t["series"]]
        try:
            chart.font.size = Pt(11.5)
            chart.font.name = self.fb
            chart.font.color.rgb = _rgb(pal["muted"])
        except Exception:
            pass
        chart.has_title = False
        chart.chart_style = None
        try:
            gf.chart_part.chart_workbook.xlsx_part = gf.chart_part.chart_workbook.xlsx_part
        except Exception:
            pass
        # Rahmen entfernen
        try:
            chart.plots[0].vary_by_categories = is_pie
        except Exception:
            pass

        show_legend = len(series) > 1 or is_pie
        chart.has_legend = show_legend
        if show_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            try:
                chart.legend.font.size = Pt(11)
                chart.legend.font.color.rgb = _rgb(pal["muted"])
            except Exception:
                pass

        plot = chart.plots[0]
        if is_pie:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.show_percentage = bool(s.get("show_percent", True))
            dl.show_value = not bool(s.get("show_percent", True))
            dl.number_format_is_linked = False
            dl.number_format = "0%" if bool(s.get("show_percent", True)) else "#,##0"
            try:
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass
            dl.font.size = Pt(11)
            dl.font.color.rgb = _rgb(pal["ink"])
            try:
                for i, pt in enumerate(chart.series[0].points):
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = _rgb(colors[i % len(colors)])
                    pt.format.line.color.rgb = _rgb(pal["bg"])
                    pt.format.line.width = Pt(1.5)
            except Exception:
                pass
        else:
            try:
                plot.gap_width = 60 if len(series) == 1 else 90
                plot.overlap = -10 if len(series) > 1 else 0
            except Exception:
                pass
            for i, ser in enumerate(chart.series):
                col = colors[i % len(colors)]
                if ctype in (XL_CHART_TYPE.LINE,):
                    ser.format.line.color.rgb = _rgb(col)
                    ser.format.line.width = Pt(2.5)
                    try:
                        ser.smooth = False
                    except Exception:
                        pass
                else:
                    ser.format.fill.solid()
                    ser.format.fill.fore_color.rgb = _rgb(col)
                    ser.format.line.fill.background()
            if s.get("show_values", len(series) == 1 and len(cats) <= 8):
                try:
                    plot.has_data_labels = True
                    dl = plot.data_labels
                    dl.font.size = Pt(11)
                    dl.font.bold = True
                    dl.font.color.rgb = _rgb(pal["ink"])
                    dl.number_format_is_linked = False
                    dl.number_format = _as_text(s.get("number_format")) or "#,##0.##"
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
                except Exception:
                    pass
            try:
                cax = chart.category_axis
                cax.has_major_gridlines = False
                cax.format.line.color.rgb = _rgb(pal["hair"])
                cax.tick_labels.font.size = Pt(11)
                cax.tick_labels.font.color.rgb = _rgb(pal["muted"])
            except Exception:
                pass
            try:
                vax = chart.value_axis
                vax.has_major_gridlines = True
                vax.major_gridlines.format.line.color.rgb = _rgb(pal["hair"])
                vax.major_gridlines.format.line.width = Pt(0.75)
                vax.format.line.fill.background()
                vax.tick_labels.font.size = Pt(11)
                vax.tick_labels.font.color.rgb = _rgb(pal["muted"])
                if s.get("number_format"):
                    vax.tick_labels.number_format_is_linked = False
                    vax.tick_labels.number_format = _as_text(s.get("number_format"))
            except Exception:
                pass
        self._footer(slide, dark, self._number)
        return slide

    def _slide_image(self, s: Dict[str, Any]):
        dark = self._is_dark("image")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        data = self._load_image(s)
        placement = str(s.get("placement") or s.get("layout") or "right").lower()
        full = placement in ("full", "cover", "ganz", "vollbild")

        if full and data:
            self._place_image_cover(slide, data, 0, 0, SW, SH)
            veil = self._shape(slide, MSO_SHAPE.RECTANGLE, 0, 3.4, SW, 4.1, fill=self.t["deep"])
            self._set_alpha(veil, 0.35)
            self._text(slide, MX, 4.55, 9.6, 1.05, _as_text(s.get("title")), size=36, bold=True,
                       color=self.t["on_deep"], font=self.fh, line=1.06, min_size=20)
            sub = _as_text(s.get("subtitle"))
            if sub:
                self._text(slide, MX, 5.75, 8.4, 0.60, sub, size=15,
                           color=self.t["on_deep_mut"], line=1.24, min_size=11)
            self._footer(slide, True, self._number)
            return slide

        img_left = placement != "left"
        panel_w = 5.85
        px = SW - panel_w if img_left else 0.0
        if data:
            self._place_image_cover(slide, data, px, 0, panel_w, SH)
        else:
            self._shape(slide, MSO_SHAPE.RECTANGLE, px, 0, panel_w, SH,
                        fill=_mix(pal["bg"], self.t["accent"], 0.16))
            self._shape(slide, MSO_SHAPE.OVAL, px + panel_w / 2 - 1.1, SH / 2 - 1.1, 2.2, 2.2,
                        fill=None, line_color=_mix(pal["bg"], self.t["accent"], 0.45), line_w=1.5)

        tx = MX if img_left else panel_w + 0.7
        tw = SW - panel_w - MX - 0.7
        kicker = _as_text(s.get("kicker"))
        title = _as_text(s.get("title"))
        sub = _as_text(s.get("subtitle"))
        bullets = _bullet_items(s.get("bullets"))[:5]

        t_size = t_h = s_size = s_h = 0.0
        if title:
            t_size, _n, t_h = self._measure(title, tw, 30, bold=True, line=1.10,
                                            max_lines=3, min_size=19)
        if sub:
            s_size, _n, s_h = self._measure(sub, tw, 15, line=1.30, max_lines=4,
                                            min_size=11)
        rows = []
        for it in bullets:
            txt = f"{it['title']} — {it['text']}" if it["title"] and it["text"] else (it["title"] or it["text"])
            bs, _n, bh = self._measure(txt, tw - 0.42, 14, line=1.26, max_lines=3,
                                       min_size=10)
            rows.append((txt, bs, bh))
        block = ((0.42 if kicker else 0.0) + (t_h + 0.26 if title else 0.0)
                 + (s_h + 0.36 if sub else 0.0)
                 + sum(bh + 0.20 for _t, _s, bh in rows))
        y = max(0.85, (SH - block) / 2)

        if kicker:
            self._text(slide, tx, y, tw, 0.26, kicker, size=10.5, bold=True,
                       color=pal["accent"], caps=True, spacing=1.4, shrink=False)
            y += 0.42
        if title:
            self._text(slide, tx, y, tw, t_h + 0.10, title, size=t_size, bold=True,
                       color=pal["ink"], font=self.fh, line=1.10, shrink=False)
            y += t_h + 0.26
        if sub:
            self._text(slide, tx, y, tw, s_h + 0.08, sub, size=s_size, color=pal["muted"],
                       line=1.30, shrink=False)
            y += s_h + 0.36
        for txt, bs, bh in rows:
            if y + bh > BODY_BOTTOM:
                break
            self._shape(slide, MSO_SHAPE.OVAL, tx + 0.02,
                        y + (bs * 1.26 / 72 - 0.105) / 2, 0.105, 0.105, fill=pal["accent"])
            self._text(slide, tx + 0.38, y, tw - 0.42, bh + 0.06, txt, size=bs,
                       color=pal["ink"], line=1.26, shrink=False)
            y += bh + 0.20
        self._footer(slide, dark, self._number)
        return slide

    def _slide_closing(self, s: Dict[str, Any]):
        dark = self._is_dark("closing")
        slide = self._new_slide(dark)
        pal = self._pal(dark)
        if dark:
            self._stage_decor(slide, 0)
        else:
            self._shape(slide, MSO_SHAPE.OVAL, 9.4, -1.5, 5.8, 5.8,
                        fill=_mix(self.t["light"], self.t["accent"], 0.10))
        title = _as_text(s.get("title")) or "Vielen Dank"
        self._text(slide, MX, 2.55, 8.4, 1.25, title, size=44, bold=True, color=pal["ink"],
                   font=self.fh, line=1.05, min_size=26)
        sub = _as_text(s.get("subtitle"))
        if sub:
            self._text(slide, MX, 3.95, 7.4, 0.85, sub, size=17, color=pal["muted"],
                       line=1.28, min_size=12)
        lines = s.get("contact") or s.get("bullets") or []
        if isinstance(lines, str):
            lines = [x for x in lines.split("\n") if x.strip()]
        y = 5.05
        for entry in (lines or [])[:4]:
            txt = _as_text(entry.get("text") if isinstance(entry, dict) else entry)
            if not txt:
                continue
            self._shape(slide, MSO_SHAPE.OVAL, MX + 0.02, y + 0.10, 0.105, 0.105,
                        fill=pal["accent"])
            self._text(slide, MX + 0.36, y, 8.0, 0.32, txt, size=13.5, color=pal["muted"],
                       shrink=False)
            y += 0.42
        self._place_logo(slide, dark)
        return slide

    # -------------------------------------------------------------- Hilfsmittel
    def _set_alpha(self, shape, alpha: float):
        """Transparenz auf eine Volltonfüllung legen (0..1 = Deckkraft)."""
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            spPr = shape.fill._xPr
            solid = spPr.find(qn("a:solidFill"))
            if solid is None:
                return
            clr = solid.find(qn("a:srgbClr"))
            if clr is None:
                return
            a = etree.SubElement(clr, qn("a:alpha"))
            a.set("val", str(int(max(0.0, min(1.0, alpha)) * 100000)))
        except Exception:
            pass

    def _load_image(self, s: Dict[str, Any]) -> Optional[bytes]:
        b64 = s.get("image_base64") or s.get("image_b64")
        if b64:
            try:
                raw = re.sub(r"^data:[^;]+;base64,", "", _as_text(b64))
                return base64.b64decode(raw)
            except Exception:
                self.warnings.append("Bild (base64) konnte nicht dekodiert werden.")
        url = _as_text(s.get("image_url"))
        if url.startswith("http"):
            try:
                import requests

                resp = requests.get(url, timeout=12, headers={"User-Agent": "OpenWebUI-PPTX/1.0"})
                resp.raise_for_status()
                if len(resp.content) > 12_000_000:
                    raise ValueError("Bild zu groß")
                return resp.content
            except Exception as exc:
                self.warnings.append(f"Bild konnte nicht geladen werden ({url[:60]}): {exc}")
        return None

    def _place_image_cover(self, slide, data: bytes, left, top, width, height):
        """Bild formatfüllend in ein Rechteck einpassen (Center-Crop)."""
        try:
            stream = io.BytesIO(data)
            pic = slide.shapes.add_picture(stream, Inches(left), Inches(top))
            nat_w, nat_h = pic.width, pic.height
            tgt_w, tgt_h = Inches(width), Inches(height)
            scale = max(tgt_w / nat_w, tgt_h / nat_h)
            new_w, new_h = int(nat_w * scale), int(nat_h * scale)
            crop_x = (new_w - tgt_w) / new_w / 2 if new_w > tgt_w else 0
            crop_y = (new_h - tgt_h) / new_h / 2 if new_h > tgt_h else 0
            pic.crop_left = crop_x
            pic.crop_right = crop_x
            pic.crop_top = crop_y
            pic.crop_bottom = crop_y
            pic.left, pic.top = Inches(left), Inches(top)
            pic.width, pic.height = tgt_w, tgt_h
            return pic
        except Exception as exc:
            self.warnings.append(f"Bild konnte nicht eingefügt werden: {exc}")
            return None

    def _place_logo(self, slide, dark: bool):
        if not self.logo_bytes:
            return
        try:
            pic = slide.shapes.add_picture(io.BytesIO(self.logo_bytes),
                                           Inches(SW - MX - 1.5), Inches(0.58))
            ratio = pic.width / pic.height if pic.height else 3
            h = Inches(0.46)
            w = int(h * ratio)
            pic.height, pic.width = h, w
            pic.left = Inches(SW - MX) - w
            pic.top = Inches(0.55)
        except Exception as exc:
            self.warnings.append(f"Logo konnte nicht platziert werden: {exc}")

    # ------------------------------------------------------------------- Build
    def build(self) -> bytes:
        renderers = {
            "title": self._slide_title, "section": self._slide_section,
            "bullets": self._slide_bullets, "agenda": self._slide_agenda,
            "cards": self._slide_cards, "stats": self._slide_stats,
            "quote": self._slide_quote, "compare": self._slide_compare,
            "table": self._slide_table, "timeline": self._slide_timeline,
            "chart": self._slide_chart, "image": self._slide_image,
            "closing": self._slide_closing,
        }
        specs = self._expand(self.slides)
        for spec in specs:
            kind = spec.get("type", "bullets")
            fn = renderers.get(kind, self._slide_bullets)
            if kind not in ("title",):
                self._number += 1
            try:
                slide = fn(spec)
            except Exception as exc:  # eine kaputte Folie darf das Deck nicht killen
                self.warnings.append(f"Folie '{_as_text(spec.get('title'))[:40]}' "
                                     f"({kind}) fehlgeschlagen: {exc}")
                slide = self._slide_bullets({"type": "bullets", "title": spec.get("title", ""),
                                             "bullets": []})
            notes = _as_text(spec.get("notes"))
            if notes:
                try:
                    slide.notes_slide.notes_text_frame.text = notes
                except Exception:
                    pass
        self._set_core_properties()
        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    def _expand(self, slides: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Lange Listen automatisch auf Folgefolien verteilen."""
        out: List[Dict[str, Any]] = []
        for s in slides:
            if s.get("type") in ("bullets", "agenda"):
                items = _bullet_items(s.get("bullets"))
                limit = 10 if s.get("type") == "agenda" else 8
                if len(items) > limit:
                    chunks = [items[i:i + limit] for i in range(0, len(items), limit)]
                    for i, chunk in enumerate(chunks):
                        clone = dict(s)
                        clone["bullets"] = chunk
                        if i:
                            clone["title"] = f"{_as_text(s.get('title'))} (Forts.)"
                            clone.pop("subtitle", None)
                            clone.pop("notes", None)
                        out.append(clone)
                    continue
            out.append(s)
        return out

    def _set_core_properties(self):
        try:
            cp = self.prs.core_properties
            cp.title = self.deck_title
            cp.subject = self.deck_subtitle
            if self.author:
                cp.author = self.author
            cp.comments = "Erstellt mit dem Open-WebUI Präsentations-Generator"
        except Exception:
            pass


# =============================================================================
#  Dateiablage in Open WebUI
# =============================================================================


def _slugify(text: str, fallback: str = "praesentation") -> str:
    trans = {"ä": "ae", "ö": "oe", "ü": "ue", "Ä": "Ae", "Ö": "Oe", "Ü": "Ue", "ß": "ss"}
    for k, v in trans.items():
        text = text.replace(k, v)
    slug = re.sub(r"[^A-Za-z0-9]+", "-", text).strip("-").lower()
    slug = re.sub(r"-{2,}", "-", slug)
    return (slug[:60] or fallback)


TOOL_VERSION = "1.3.0"

PPTX_MIME = ("application/vnd.openxmlformats-officedocument."
             "presentationml.presentation")


async def _resolve(value: Any) -> Any:
    """Wartet auf Rückgaben ab, die je nach Open-WebUI-Version async sind.

    Ab Open WebUI 0.9.0 sind die Modell- und Storage-Methoden Coroutinen. Ein nicht
    abgewartetes Ergebnis ist ein wahrheitswertiges Coroutine-Objekt – der Aufruf sieht
    dann erfolgreich aus, obwohl nichts passiert ist.
    """
    if inspect.isawaitable(value):
        return await value
    return value


def _upload_dir() -> str:
    """Upload-Verzeichnis der laufenden Instanz ermitteln (versionsabhängig)."""
    for module, attr in (("open_webui.config", "UPLOAD_DIR"),
                         ("open_webui.env", "UPLOAD_DIR"),
                         ("open_webui.config", "DATA_DIR"),
                         ("open_webui.env", "DATA_DIR")):
        try:
            mod = __import__(module, fromlist=[attr])
            value = str(getattr(mod, attr) or "")
            if value:
                return value if attr == "UPLOAD_DIR" else os.path.join(value, "uploads")
        except Exception:
            continue
    return os.path.join(os.environ.get("DATA_DIR", "/app/backend/data"), "uploads")


def _file_form_fields() -> set:
    """Felder, die das FileForm dieser Open-WebUI-Version kennt."""
    from open_webui.models.files import FileForm

    fields = getattr(FileForm, "model_fields", None) or getattr(FileForm, "__fields__", {})
    return set(fields or ())


def _local_path(path: str) -> Optional[str]:
    """Lokaler Pfad zu einem Storage-Pfad; None bei Objektspeicher (S3/GCS/Azure)."""
    if not path:
        return None
    if "://" in path and not path.startswith("file://"):
        return None
    return path[7:] if path.startswith("file://") else path


async def _verify_stored_file(file_id: str, user_id: str,
                              expected_size: Optional[int] = None) -> Tuple[bool, str]:
    """Prüft, ob /api/v1/files/<id>/content die Datei wirklich ausliefern kann.

    Geprüft wird gegen genau das, was der Endpunkt tut: Eintrag lesen, Besitzer
    vergleichen, Pfad auflösen, Datei öffnen. Alles wird abgewartet – sonst prüft man
    Coroutine-Objekte statt Daten und bekommt ein grünes Ergebnis für einen leeren Vorgang.
    """
    from open_webui.models.files import Files

    record = await _resolve(Files.get_file_by_id(file_id))
    if record is None:
        return False, "kein Datenbank-Eintrag angelegt"
    if not getattr(record, "id", None) and not isinstance(record, dict):
        return False, (f"Datenbank lieferte kein Datei-Objekt, sondern "
                       f"{type(record).__name__}")
    if isinstance(record, dict):
        record = types.SimpleNamespace(**record)

    owner = getattr(record, "user_id", None)
    if owner and user_id and owner != user_id:
        return False, "Eintrag gehört einem anderen Benutzerkonto"

    path = str(getattr(record, "path", "") or "")
    if path:
        local: Optional[str] = None
        try:
            from open_webui.storage.provider import Storage

            resolved = await _resolve(Storage.get_file(path))
            local = str(resolved) if resolved else None
        except Exception:
            local = None
        if local is None:
            local = _local_path(path)
        if local is None:
            return True, ""  # Objektspeicher ohne lokale Entsprechung
        if not os.path.isfile(local):
            return False, f"keine Datei unter {local}"
        if expected_size and os.path.getsize(local) != expected_size:
            return False, (f"Datei unter {local} hat {os.path.getsize(local)} statt "
                           f"{expected_size} Bytes")
        return True, ""

    # Ältere Versionen ohne path-Spalte: der Endpunkt sucht UPLOAD_DIR/<filename>
    name = str(getattr(record, "filename", "") or "")
    if not name:
        return False, "Eintrag ohne Pfad und ohne Dateinamen"
    legacy = os.path.join(_upload_dir(), name)
    if os.path.isfile(legacy):
        return True, ""
    payload = getattr(record, "data", None) or {}
    if isinstance(payload, dict) and payload.get("content"):
        return True, ""
    return False, f"keine Datei unter {legacy}"


async def _insert_file_record(file_id: str, data: bytes, storage_path: Optional[str],
                              display_name: str, stored_name: str, user_id: str):
    """Legt den Datei-Eintrag passend zum FileForm der installierten Version an."""
    from open_webui.models.files import FileForm, Files

    supported = _file_form_fields()
    has_path = "path" in supported
    fields = {
        "id": file_id,
        # Ohne path-Spalte liest der Download-Endpunkt UPLOAD_DIR/<filename> –
        # dann muss dort der gespeicherte (präfixierte) Name stehen, nicht der schöne.
        "filename": display_name if has_path else stored_name,
        "path": storage_path,
        "data": {},
        "meta": {"name": display_name, "content_type": PPTX_MIME, "size": len(data)},
        "hash": hashlib.sha256(data).hexdigest(),
        "access_control": None,
    }
    form = FileForm(**{k: v for k, v in fields.items() if k in supported})
    record = await _resolve(Files.insert_new_file(user_id, form))
    # insert_new_file fängt Fehler intern ab und liefert dann None statt zu werfen.
    if record is None:
        raise RuntimeError("insert_new_file() lieferte None – der Datenbank-Eintrag "
                           "wurde nicht angelegt")
    return record


async def _store_via_router(data: bytes, filename: str, user_id: str, request) -> str:
    """Weg 1: die offizielle Upload-Route aufrufen – passt sich der Version selbst an."""
    from fastapi import BackgroundTasks, UploadFile
    from open_webui.models.users import Users
    from starlette.datastructures import Headers

    routers = __import__("open_webui.routers.files", fromlist=["upload_file_handler"])
    # 0.9.0 heißt die Funktion upload_file_handler, davor upload_file.
    upload_route = (getattr(routers, "upload_file_handler", None)
                    or getattr(routers, "upload_file", None))
    if upload_route is None:
        raise ImportError("weder upload_file_handler noch upload_file vorhanden")

    user = await _resolve(Users.get_user_by_id(user_id))
    if user is None:
        raise RuntimeError("Benutzerobjekt nicht gefunden")
    headers = Headers({"content-type": PPTX_MIME})
    try:
        upload = UploadFile(file=io.BytesIO(data), filename=filename, headers=headers)
    except TypeError:
        upload = UploadFile(filename=filename, file=io.BytesIO(data))

    kwargs: Dict[str, Any] = {}
    for name in inspect.signature(upload_route).parameters:
        if name == "request":
            if request is None:
                raise RuntimeError("kein __request__ verfügbar")
            kwargs[name] = request
        elif "background" in name and "process" not in name:
            kwargs[name] = BackgroundTasks()
        elif name == "file":
            kwargs[name] = upload
        elif name == "user":
            kwargs[name] = user
        elif name in ("process", "process_in_background"):
            kwargs[name] = False  # keine RAG-Indizierung für eine PPTX
        elif name in ("file_metadata", "metadata"):
            kwargs[name] = {"generated_by": "praesentations_generator"}
        elif name == "internal":
            kwargs[name] = True

    result = await _resolve(upload_route(**kwargs))
    file_id = getattr(result, "id", None)
    if file_id is None and isinstance(result, dict):
        file_id = result.get("id")
    if not file_id:
        raise RuntimeError(f"Upload-Route lieferte keine Datei-ID "
                           f"({type(result).__name__})")
    return str(file_id)


async def _store_via_storage_provider(data: bytes, filename: str, user_id: str) -> str:
    """Weg 2: Storage-Provider (unterstützt auch S3/GCS) plus eigener DB-Eintrag."""
    from open_webui.storage.provider import Storage

    file_id = str(uuid.uuid4())
    stored_name = f"{file_id}_{filename}"
    buf = io.BytesIO(data)
    buf.name = stored_name
    try:
        result = await _resolve(Storage.upload_file(buf, stored_name, {}))
    except TypeError:
        buf.seek(0)
        result = await _resolve(Storage.upload_file(buf, stored_name))
    if isinstance(result, (tuple, list)) and len(result) >= 2:
        path = result[1]
    else:
        path = result
    if not isinstance(path, (str, os.PathLike)):
        raise RuntimeError(f"Storage lieferte keinen Pfad, sondern "
                           f"{type(path).__name__}")
    await _insert_file_record(file_id, data, str(path), filename, stored_name, user_id)
    return file_id


async def _store_via_disk(data: bytes, filename: str, user_id: str) -> str:
    """Weg 3: direkt ins Upload-Verzeichnis schreiben plus DB-Eintrag."""
    file_id = str(uuid.uuid4())
    stored_name = f"{file_id}_{filename}"
    upload_dir = _upload_dir()
    os.makedirs(upload_dir, exist_ok=True)
    path = os.path.join(upload_dir, stored_name)
    with open(path, "wb") as fh:
        fh.write(data)
    await _insert_file_record(file_id, data, path, filename, stored_name, user_id)
    return file_id


async def store_in_open_webui(data: bytes, filename: str, user_id: str,
                              request: Any = None) -> Dict[str, Any]:
    """Legt die Datei ab und gibt erst dann einen Link zurück, wenn er auch trägt.

    Drei Wege werden der Reihe nach versucht; jeder wird danach überprüft. Nur ein
    nachweislich abrufbarer Eintrag führt zu einem Download-Link.
    """
    if not user_id:
        return {"error": "Keine Benutzer-ID verfügbar – die Datei kann keinem Konto "
                         "zugeordnet werden."}
    try:
        import open_webui.models.files  # noqa: F401
    except Exception as exc:
        return {"error": f"Open-WebUI-Dateimodell nicht verfügbar: {exc}"}

    problems: List[str] = []
    strategies = (
        ("Upload-Route", lambda: _store_via_router(data, filename, user_id, request)),
        ("Storage-Provider", lambda: _store_via_storage_provider(data, filename, user_id)),
        ("Upload-Verzeichnis", lambda: _store_via_disk(data, filename, user_id)),
    )
    for label, run in strategies:
        try:
            file_id = await _resolve(run())
        except Exception as exc:
            problems.append(f"{label} – {type(exc).__name__}: {exc}")
            continue
        try:
            ok, why = await _verify_stored_file(file_id, user_id, len(data))
        except Exception as exc:
            problems.append(f"{label} – Prüfung fehlgeschlagen: "
                            f"{type(exc).__name__}: {exc}")
            continue
        if ok:
            return {"id": file_id,
                    "url": f"/api/v1/files/{file_id}/content?attachment=true",
                    "via": label, "notes": problems}
        problems.append(f"{label} – {why}")
    return {"error": " | ".join(problems) or "unbekannter Fehler"}


# =============================================================================
#  Open-WebUI-Tool
# =============================================================================


class Tools:
    class Valves(BaseModel):
        default_theme: str = Field(
            default="midnight",
            description=("Standard-Theme, wenn das Modell keines wählt: midnight, graphite, "
                         "teal, ocean, forest, terracotta, coral, berry, cherry"),
        )
        mode: str = Field(
            default="auto",
            description=("Helligkeit: 'auto' (dunkle Titel-/Kapitelfolien, helle Inhalte), "
                         "'light' (durchgehend hell) oder 'dark' (durchgehend dunkel)"),
        )
        brand_accent: str = Field(
            default="",
            description="Eigene Akzentfarbe als Hex (z. B. #C8102E). Leer = Theme-Farbe.",
        )
        brand_accent2: str = Field(
            default="", description="Zweite Akzentfarbe als Hex (optional).",
        )
        brand_deep: str = Field(
            default="",
            description="Eigene dunkle Bühnenfarbe für Titel-/Kapitelfolien als Hex (optional).",
        )
        font_heading: str = Field(
            default="Cambria",
            description="Schrift für Überschriften. Sichere Office-Schriften: Cambria, "
                        "Calibri, Arial, Century Schoolbook, Bookman Old Style.",
        )
        font_body: str = Field(
            default="Calibri", description="Schrift für Fließtext.",
        )
        footer_text: str = Field(
            default="",
            description="Fester Fußzeilentext (z. B. Firmenname). Leer = Titel der Präsentation.",
        )
        show_footer: bool = Field(default=True, description="Fußzeile anzeigen.")
        show_slide_numbers: bool = Field(default=True, description="Foliennummern anzeigen.")
        logo_url: str = Field(
            default="",
            description="URL zu einem Logo (PNG/JPG) für Titel- und Abschlussfolie.",
        )
        author_default: str = Field(
            default="", description="Standard-Autor in der Fußzeile der Titelfolie.",
        )
        max_slides: int = Field(default=40, description="Maximale Anzahl Folien pro Datei.")
        base_url: str = Field(
            default="",
            description="Optionale externe Basis-URL (z. B. https://chat.firma.at) für "
                        "absolute Download-Links.",
        )
        attach_to_message: bool = Field(
            default=True,
            description="Datei zusätzlich als Anhang an die Chat-Nachricht hängen.",
        )

    class UserValves(BaseModel):
        theme: str = Field(
            default="",
            description="Persönliches Lieblings-Theme (überschreibt den Admin-Standard).",
        )
        author: str = Field(default="", description="Eigener Name für die Titelfolie.")
        language_hint: str = Field(
            default="de", description="Sprache der Präsentation (nur informativ)."
        )

    def __init__(self):
        self.valves = self.Valves()
        self.user_valves = self.UserValves()
        self.citation = False  # eigene Quellenangabe, keine automatische

    # ------------------------------------------------------------------ Themes
    def list_themes(self) -> str:
        """
        Listet alle verfügbaren Design-Themes für Präsentationen mit kurzer Beschreibung auf.
        Nutze das, wenn der Nutzer fragt, welche Designs es gibt, oder wenn du ein zum Thema
        passendes Design auswählen möchtest.
        :return: Liste der Theme-Namen mit Beschreibung.
        """
        lines = ["Verfügbare Themes für create_presentation(theme=...):", ""]
        for key, t in THEMES.items():
            lines.append(f"- **{key}** – {t['label']}")
        lines += [
            "",
            "Modus (mode): 'auto' = dunkle Titel-/Kapitelfolien + helle Inhaltsfolien, "
            "'light' = durchgehend hell, 'dark' = durchgehend dunkel (Premium-Look).",
        ]
        return "\n".join(lines)

    # -------------------------------------------------------------- Diagnose
    async def check_setup(
        self,
        __user__: Optional[dict] = None,
        __request__: Optional[Any] = None,
    ) -> str:
        """
        Prüft, ob der Präsentations-Generator in dieser Open-WebUI-Installation sauber
        arbeitet: Bibliotheken, Datei-Schnittstelle, Upload-Verzeichnis, Schreibrechte.
        Legt dazu testweise eine kleine Datei an, prüft sie und löscht sie wieder.
        Nutze das, wenn ein Download-Link nicht funktioniert oder ins Leere führt.
        :return: Prüfbericht im Klartext.
        """
        lines: List[str] = [f"**Setup-Prüfung Präsentations-Generator {TOOL_VERSION}**", ""]
        ok = True

        try:
            import pptx as _pptx

            lines.append(f"✅ python-pptx {getattr(_pptx, '__version__', '?')}")
        except Exception as exc:
            ok = False
            lines.append(f"❌ python-pptx fehlt: {exc}")

        try:
            fields = _file_form_fields()
            has_path = "path" in fields
            lines.append(f"✅ FileForm-Felder: {', '.join(sorted(fields))}")
            if has_path:
                lines.append("✅ path-Spalte vorhanden")
            else:
                lines.append("ℹ️ Keine path-Spalte – ältere Open-WebUI-Version, "
                             "die Ablage läuft im Kompatibilitätsmodus")
        except Exception as exc:
            ok = False
            lines.append(f"❌ Datei-Schnittstelle nicht erreichbar: {exc}")

        upload_dir = _upload_dir()
        if os.path.isdir(upload_dir) and os.access(upload_dir, os.W_OK):
            lines.append(f"✅ Upload-Verzeichnis beschreibbar: `{upload_dir}`")
        else:
            lines.append(f"⚠️ Upload-Verzeichnis nicht beschreibbar oder fehlt: "
                         f"`{upload_dir}`")

        try:
            from open_webui.storage.provider import Storage

            lines.append(f"✅ Storage-Provider: {type(Storage).__name__}")
        except Exception as exc:
            lines.append(f"ℹ️ Kein Storage-Provider importierbar ({exc}) – "
                         f"es wird direkt ins Upload-Verzeichnis geschrieben.")

        user_id = (__user__ or {}).get("id") or (__user__ or {}).get("user_id") or ""
        if user_id:
            lines.append(f"✅ Benutzer-ID vorhanden: `{user_id}`")
        else:
            ok = False
            lines.append("❌ Keine Benutzer-ID – Dateien lassen sich keinem Konto zuordnen.")

        # Echter Durchlauf mit einer Mini-Datei
        if user_id:
            try:
                probe = Presentation()
                probe.slide_width, probe.slide_height = Inches(SW), Inches(SH)
                probe.slides.add_slide(probe.slide_layouts[6])
                buf = io.BytesIO()
                probe.save(buf)
                result = await store_in_open_webui(
                    buf.getvalue(), "setup-test.pptx", user_id, __request__)
                if result.get("error"):
                    ok = False
                    lines.append(f"❌ Testablage fehlgeschlagen: {result['error']}")
                else:
                    for note in result.get("notes") or []:
                        lines.append(f"ℹ️ Nicht genutzt – {note}")
                    lines.append(f"✅ Testablage erfolgreich über: {result['via']}")
                    lines.append(f"   Download-Pfad wäre: `{result['url']}`")
                    await self._remove_test_file(result["id"])
            except Exception as exc:
                ok = False
                lines.append(f"❌ Testablage abgebrochen: {type(exc).__name__}: {exc}")

        lines.append("")
        lines.append("**Ergebnis:** " + ("Alles bereit – Präsentationen können erzeugt "
                                         "und heruntergeladen werden."
                                         if ok else
                                         "Es gibt ein Problem, siehe die mit ❌ markierten "
                                         "Punkte."))
        return "\n".join(lines)

    async def check_link(self, link: str, __user__: Optional[dict] = None) -> str:
        """
        Untersucht, warum ein bestimmter Download-Link einer Präsentation nicht
        funktioniert. Übergib den Link oder die Datei-ID aus der Fehlermeldung.
        Nutze das, wenn ein Link „We could not find what you're looking for" liefert –
        statt die Präsentation einfach neu zu erzeugen.
        :param link: Der nicht funktionierende Download-Link oder die Datei-ID.
        :return: Befund und konkrete Empfehlung im Klartext.
        """
        match = re.search(
            r"[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-"
            r"[0-9a-fA-F]{4}-[0-9a-fA-F]{12}", link or "")
        if not match:
            return ("❌ In der Eingabe steckt keine Datei-ID. Bitte den vollständigen Link "
                    "übergeben, etwa `/api/v1/files/…/content`.")
        file_id = match.group(0)
        user_id = (__user__ or {}).get("id") or (__user__ or {}).get("user_id") or ""

        try:
            from open_webui.models.files import Files
        except Exception as exc:
            return f"❌ Datei-Schnittstelle nicht erreichbar: {exc}"

        record = await _resolve(Files.get_file_by_id(file_id))
        out = [f"**Prüfung des Links** (Datei-ID `{file_id}`)", ""]

        if record is None:
            out += [
                "❌ Zu dieser ID gibt es **keinen Eintrag** in der Datenbank.",
                "",
                "Das ist das Fehlerbild von Tool-Version 1.0.0: Sie hat einen Link "
                "ausgegeben, ohne dass der Eintrag angelegt wurde. Solche Links bleiben "
                "dauerhaft tot, auch nach dem Update – die Präsentation muss neu erzeugt "
                "werden.",
                "",
                "**Empfehlung:** Die Präsentation noch einmal erstellen lassen. Wenn der "
                "neue Link ebenfalls nicht funktioniert, diesen neuen Link hier erneut "
                "prüfen – dann liegt eine andere Ursache vor.",
            ]
            out += await self._recent_files(user_id)
            return "\n".join(out)

        owner = getattr(record, "user_id", None)
        meta = getattr(record, "meta", None) or {}
        out.append(f"✅ Eintrag vorhanden: `{getattr(record, 'filename', '?')}`"
                   f" ({meta.get('size', '?')} Bytes)")
        if owner and user_id and owner != user_id:
            out += ["❌ Der Eintrag gehört einem **anderen Benutzerkonto** "
                    f"(`{owner}`), du bist `{user_id}`.", "",
                    "**Empfehlung:** Die Präsentation in dem Konto erzeugen, in dem du sie "
                    "herunterladen willst – oder das Deck neu erstellen lassen."]
            return "\n".join(out)
        out.append("✅ Der Eintrag gehört deinem Konto")

        path = getattr(record, "path", None)
        if path:
            local = _local_path(str(path))
            if local is None:
                out.append(f"ℹ️ Objektspeicher-Pfad: `{path}` – Existenz nicht lokal prüfbar")
            elif os.path.exists(local):
                out.append(f"✅ Datei liegt unter `{local}` "
                           f"({os.path.getsize(local)} Bytes)")
            else:
                out += [f"❌ Unter `{local}` liegt **keine Datei**.", "",
                        "**Empfehlung:** Das passiert, wenn das Upload-Verzeichnis nicht "
                        "dauerhaft gespeichert wird (Container ohne Volume) oder wenn "
                        "mehrere Open-WebUI-Instanzen mit lokalem Speicher hinter einem "
                        "Load Balancer laufen. Volume prüfen bzw. gemeinsamen Speicher "
                        "(S3) einrichten."]
                return "\n".join(out)
        else:
            legacy = os.path.join(_upload_dir(), str(getattr(record, "filename", "")))
            state = "gefunden" if os.path.exists(legacy) else "**nicht gefunden**"
            out.append(f"{'✅' if os.path.exists(legacy) else '❌'} "
                       f"Kompatibilitätsmodus: Datei unter `{legacy}` {state}")

        out += ["",
                "**Ergebnis:** Eintrag und Datei sind vollständig – der Download über "
                f"`/api/v1/files/{file_id}/content` müsste funktionieren.",
                "",
                "Wenn er es trotzdem nicht tut, liegt es außerhalb des Tools: den Link im "
                "selben, angemeldeten Browser öffnen (nicht in einem privaten Fenster oder "
                "einem anderen Konto), und bei einem Reverse Proxy prüfen, ob "
                "`/api/v1/files/…` durchgereicht wird."]
        return "\n".join(out)

    @staticmethod
    async def _recent_files(user_id: str) -> List[str]:
        """Listet die letzten Präsentationen des Kontos – hilft beim Einordnen."""
        if not user_id:
            return []
        try:
            from open_webui.models.files import Files

            getter = getattr(Files, "get_files_by_user_id", None)
            records = await _resolve(getter(user_id)) if getter else []
        except Exception:
            return []
        decks = [r for r in (records or [])
                 if str(getattr(r, "filename", "")).lower().endswith(".pptx")][-3:]
        if not decks:
            return ["", "In diesem Konto ist bislang keine Präsentation gespeichert."]
        lines = ["", "**Zuletzt gespeicherte Präsentationen dieses Kontos:**"]
        for r in reversed(decks):
            lines.append(f"- `{getattr(r, 'filename', '?')}` → "
                         f"`/api/v1/files/{getattr(r, 'id', '?')}/content`")
        return lines

    @staticmethod
    async def _remove_test_file(file_id: str) -> None:
        """Testdatei nach der Prüfung wieder entfernen."""
        try:
            from open_webui.models.files import Files

            record = await _resolve(Files.get_file_by_id(file_id))
            path = _local_path(str(getattr(record, "path", "") or "")) if record else None
            await _resolve(Files.delete_file_by_id(file_id))
            if path and os.path.isfile(path):
                os.remove(path)
        except Exception:
            pass

    # ------------------------------------------------------------------- Build
    async def create_presentation(
        self,
        title: str,
        slides: str,
        subtitle: str = "",
        theme: str = "auto",
        mode: str = "auto",
        author: str = "",
        __user__: Optional[dict] = None,
        __event_emitter__: Optional[Callable[[dict], Any]] = None,
        __request__: Optional[Any] = None,
    ) -> str:
        """
        Erzeugt eine professionell gestaltete PowerPoint-Präsentation (16:9) und gibt einen
        Download-Link zurück. Schreibe den Folieninhalt selbst und übergib ihn als JSON.

        WICHTIG: Verfasse aussagekräftige, konkrete Inhalte (keine Platzhalter), variiere die
        Folientypen (nicht nur 'bullets'), halte Titel unter 60 Zeichen, Bullet-Texte unter
        140 Zeichen und plane 8-14 Folien für einen normalen Vortrag. Beginne mit 'title',
        ende mit 'closing'.

        Der Parameter `slides` ist ein JSON-Array. Jedes Objekt hat "type" und passende Felder.
        Optional überall: "notes" (Sprechernotizen), "kicker" (kleines Label über dem Titel).

        Verfügbare Folientypen:
        - {"type":"title","title":"...","subtitle":"...","kicker":"..."}
        - {"type":"agenda","title":"Agenda","bullets":["Punkt 1","Punkt 2"]}
        - {"type":"section","title":"Kapitelname","subtitle":"...","index":1}
        - {"type":"bullets","title":"...","subtitle":"...","bullets":[{"title":"Kurzlabel","text":"Erklärung"},"einfacher Punkt"]}
        - {"type":"cards","title":"...","cards":[{"title":"...","text":"..."}]}  (2-6 Karten)
        - {"type":"stats","title":"...","stats":[{"value":"87 %","label":"...","note":"..."}]}  (2-4 Kennzahlen)
        - {"type":"compare","title":"...","columns":[{"title":"Heute","bullets":[...]},{"title":"Morgen","bullets":[...],"highlight":true}]}
        - {"type":"timeline","title":"...","steps":[{"label":"Q1","title":"...","text":"..."}]}  (3-5 Schritte)
        - {"type":"table","title":"...","columns":["A","B"],"rows":[["a1","b1"],["a2","b2"]]}  (max. 9 Zeilen)
        - {"type":"chart","title":"...","chart_type":"column|bar|line|pie|doughnut|stacked","categories":["2023","2024"],"series":[{"name":"Umsatz","values":[12,18]}],"takeaway":"Kernaussage in einem Satz"}
        - {"type":"quote","quote":"...","attribution":"Name, Rolle"}
        - {"type":"image","title":"...","subtitle":"...","image_url":"https://...","placement":"right|left|full","bullets":[...]}
        - {"type":"closing","title":"Vielen Dank","subtitle":"...","contact":["mail@firma.at","www.firma.at"]}

        :param title: Titel der Präsentation (erscheint auf der Titelfolie und in der Fußzeile).
        :param slides: JSON-Array mit den Folien-Objekten (siehe Beschreibung oben).
        :param subtitle: Untertitel für die Titelfolie.
        :param theme: Design-Theme: midnight, graphite, teal, ocean, forest, terracotta, coral, berry, cherry. 'auto' nutzt die Voreinstellung.
        :param mode: 'auto', 'light' oder 'dark'.
        :param author: Name des Vortragenden für die Titelfolie.
        :return: Bestätigung mit Download-Link zur erzeugten .pptx-Datei.
        """
        emit = __event_emitter__

        async def status(text: str, done: bool = False):
            if emit:
                try:
                    await emit({"type": "status",
                                "data": {"description": text, "done": done}})
                except Exception:
                    pass

        await status("Folien werden gestaltet …")

        # ---- Eingaben normalisieren
        try:
            spec = parse_slides(slides)
        except Exception as exc:
            await status("Fehler in den Foliendaten", True)
            return (f"❌ Die Foliendaten konnten nicht gelesen werden: {exc}\n\n"
                    "Bitte 'slides' als JSON-Array übergeben, z. B.: "
                    '[{"type":"title","title":"Titel"},{"type":"bullets","title":"Thema",'
                    '"bullets":["Punkt A","Punkt B"]}]')
        if not spec:
            await status("Keine Folien übergeben", True)
            return "❌ Es wurden keine Folien übergeben. Bitte 'slides' mit mindestens einem Objekt füllen."

        v = self.valves
        uv = self.UserValves()
        if __user__ and isinstance(__user__.get("valves"), BaseModel):
            uv = __user__["valves"]
        elif __user__ and isinstance(__user__.get("valves"), dict):
            try:
                uv = self.UserValves(**__user__["valves"])
            except Exception:
                pass

        spec = spec[: max(1, int(v.max_slides or 40))]
        if not any(s.get("type") == "title" for s in spec[:1]):
            spec.insert(0, {"type": "title", "title": _as_text(title),
                            "subtitle": _as_text(subtitle)})

        chosen_theme = theme if theme and theme.lower() not in ("auto", "", "default") else (
            uv.theme or v.default_theme)
        theme_key = resolve_theme(chosen_theme, v.default_theme)
        deck_mode = (mode or "auto").lower()
        if deck_mode not in ("auto", "light", "dark"):
            deck_mode = v.mode if v.mode in ("auto", "light", "dark") else "auto"

        logo = self._fetch_logo(v.logo_url)

        designer = DeckDesigner(
            title=title,
            slides=spec,
            theme=theme_key,
            subtitle=subtitle,
            author=_as_text(author) or uv.author or v.author_default,
            date_text=datetime.now().strftime("%d.%m.%Y"),
            mode=deck_mode,
            font_heading=v.font_heading,
            font_body=v.font_body,
            footer_text=v.footer_text,
            show_footer=v.show_footer,
            show_slide_numbers=v.show_slide_numbers,
            brand={"accent": v.brand_accent, "accent2": v.brand_accent2, "deep": v.brand_deep},
            logo_bytes=logo,
        )

        try:
            data = await asyncio.to_thread(designer.build)
        except Exception as exc:
            await status("Erstellung fehlgeschlagen", True)
            return f"❌ Die Präsentation konnte nicht erzeugt werden: {exc}"

        n_slides = len(designer.prs.slides._sldIdLst)
        filename = f"{_slugify(_as_text(title))}.pptx"
        await status("Datei wird bereitgestellt …")

        user_id = (__user__ or {}).get("id") or (__user__ or {}).get("user_id") or ""
        stored = await store_in_open_webui(data, filename, user_id, __request__)

        if stored.get("error"):
            await status("Datei konnte nicht abgelegt werden", True)
            return (f"❌ Die Präsentation wurde erstellt, konnte aber nicht in Open WebUI "
                    f"abgelegt werden – deshalb gibt es keinen Download-Link.\n\n"
                    f"Details: {stored['error']}\n\n"
                    f"Sag dem Nutzer, er soll die Funktion `check_setup` dieses Tools "
                    f"aufrufen; sie prüft Datei-Schnittstelle, Upload-Verzeichnis und "
                    f"Schreibrechte und benennt die Ursache.")

        url = stored["url"]
        if v.base_url:
            url = v.base_url.rstrip("/") + url

        theme_label = THEMES[theme_key]["label"].split("–")[0].strip()
        msg = (f"\n\n📊 **[{filename}]({url})** &nbsp;·&nbsp; {n_slides} Folien &nbsp;·&nbsp; "
               f"Design: {theme_label}\n")
        if designer.warnings:
            msg += "\n> Hinweise: " + "; ".join(designer.warnings[:4]) + "\n"

        if emit:
            try:
                await emit({"type": "message", "data": {"content": msg}})
            except Exception:
                pass
            if v.attach_to_message:
                try:
                    await emit({
                        "type": "files",
                        "data": {"files": [{
                            "type": "file",
                            "id": stored["id"],
                            "name": filename,
                            "url": f"/api/v1/files/{stored['id']}",
                            "collection_name": "",
                            "status": "uploaded",
                        }]},
                    })
                except Exception:
                    pass
        await status(f"Präsentation fertig · {n_slides} Folien", True)

        return (f"✅ Die Präsentation „{_as_text(title)}“ wurde mit {n_slides} Folien im Design "
                f"„{theme_label}“ erstellt. Der Download-Link steht bereits in der Antwort – "
                f"gib ihn nicht erneut aus, fasse stattdessen kurz den Aufbau des Decks zusammen."
                + (f" Hinweise: {'; '.join(designer.warnings[:4])}" if designer.warnings else ""))

    # ------------------------------------------------------------------ Helfer
    def _fetch_logo(self, url: str) -> Optional[bytes]:
        url = (url or "").strip()
        if not url.startswith("http"):
            return None
        try:
            import requests

            resp = requests.get(url, timeout=8)
            resp.raise_for_status()
            if len(resp.content) > 4_000_000:
                return None
            return resp.content
        except Exception:
            return None
